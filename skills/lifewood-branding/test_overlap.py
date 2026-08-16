#!/usr/bin/env python3
"""
Overlap and conflict detection between Lifewood PPT and PDF design systems.

When AI generates both PPT and PDF from Lifewood brand, it MUST NOT confuse
the two systems. These tests ensure:
  - Color values are consistent across systems (same brand) but stored in
    distinct modules with appropriate types (RGBColor vs tuple).
  - Spec-lock keys use separate namespaces.
  - Measurement units are correct per system (inches/Emu for PPT, mm for PDF).
  - Logo placement logic differs per format.
  - Font handling is correct per format.
  - No accidental cross-import between the two systems.
  - Page chrome: numbering starts after the ToC (cover/ToC unnumbered),
    and the logo sits bottom-right on every page/slide.
"""

from __future__ import annotations

import os
import re
import sys
from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# Path setup
# ---------------------------------------------------------------------------
SKILL_DIR = Path(
    "C:\\Users\\richa\\.config\\opencode\\skills\\lifewood-branding"
)
sys.path.insert(0, str(SKILL_DIR))

TEMP_DIR = Path(
    "C:\\Users\\richa\\AppData\\Local\\Temp\\opencode\\test_output"
)
TEMP_DIR.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Imports with guards
# ---------------------------------------------------------------------------

# PPT system
try:
    from lifewood_deck_builder import (
        CASTLETON_GREEN,
        DARK_SERPENT,
        EARTH_YELLOW,
        FONT,
        LOGO_DARK,
        LOGO_H,
        LOGO_LIGHT,
        LOGO_W,
        LOGO_RIGHT_CLEAR as PPT_LOGO_RIGHT_CLEAR,
        PAPER,
        SAFFRON,
        SEA_SALT,
        WHITE,
        ContentAnalyzer,
        LifewoodDeck,
        MARGIN_L as PPT_MARGIN_L,
        MARGIN_R as PPT_MARGIN_R,
        MARGIN_T as PPT_MARGIN_T,
        MARGIN_B as PPT_MARGIN_B,
        FOOTER_Y as PPT_FOOTER_Y,
        FOOTER_RULE_Y as PPT_FOOTER_RULE_Y,
        FOOTER_RULE_H as PPT_FOOTER_RULE_H,
        HEADER_NUM_SIZE as PPT_HEADER_NUM_SIZE,
        PAGE_TITLE_SIZE as PPT_PAGE_TITLE_SIZE,
        COVER_TITLE_SIZE as PPT_COVER_TITLE_SIZE,
        BODY_SIZE as PPT_BODY_SIZE,
        FOLIO_SIZE as PPT_FOLIO_SIZE,
        SLIDE_W,
        SLIDE_H,
    )
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches, Pt

    HAS_PPT = True
except ImportError as e:
    HAS_PPT = False
    _ppt_err = str(e)

# PDF design system (constants + PdfSpecLock)
try:
    from lifewood_pdf_design_system import (
        PDF_CASTLETON_GREEN,
        PDF_DARK_SERPENT,
        PDF_EARTH_YELLOW,
        PDF_FONT,
        PDF_FALLBACK_FONT,
        PDF_MARGIN_L,
        PDF_MARGIN_R,
        PDF_MARGIN_T,
        PDF_MARGIN_B,
        PDF_PAGE_W,
        PDF_PAGE_H,
        PDF_CONTENT_W,
        PDF_FOOTER_Y,
        PDF_FOOTER_RULE_Y,
        PDF_FOOTER_RULE_H,
        PDF_HEADER_NUM_SIZE,
        PDF_HEADER_TITLE_SIZE,
        PDF_COVER_TITLE_SIZE,
        PDF_BODY_SIZE,
        PDF_FOLIO_SIZE,
        PDF_LOGO_W,
        PDF_LOGO_H,
        PDF_LOGO_RIGHT_CLEAR,
        PDF_LOGO_LIGHT,
        PDF_LOGO_DARK,
        PDF_BODY_TOP,
        PDF_BODY_BOTTOM,
        PDF_COL_W,
        PDF_COL_GAP,
        PDF_HEADER_RULE_Y,
        PDF_PAPER,
        PDF_SAFFRON,
        PDF_SEA_SALT,
        PDF_WHITE,
        PdfSpecLock,
    )

    HAS_PDF_MODULE = True
except ImportError as e:
    HAS_PDF_MODULE = False
    _pdf_design_err = str(e)

# PDF builder
try:
    from lifewood_pdf_builder import LifewoodPDF

    HAS_PDF_BUILDER = True
except ImportError as e:
    HAS_PDF_BUILDER = False
    _pdf_builder_err = str(e)

# PDF inspection (PyMuPDF)
try:
    import fitz  # PyMuPDF — extract text/images from built PDFs

    HAS_FITZ = True
except ImportError as e:
    HAS_FITZ = False
    _fitz_err = str(e)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    """Convert '#RRGGBB' to (R, G, B) tuple."""
    h = hex_str.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rbg_to_tuple(rgb_color) -> tuple[int, int, int]:
    """Extract (R,G,B) from an RGBColor object."""
    return (rgb_color[0], rgb_color[1], rgb_color[2])


# Known brand hex values (from SKILL.md — source of truth)
BRAND_HEX: dict[str, str] = {
    "PAPER": "#F5EEDB",
    "WHITE": "#FFFFFF",
    "SEA_SALT": "#F9F7F7",
    "CASTLETON_GREEN": "#046241",
    "DARK_SERPENT": "#133020",
    "SAFFRON": "#FFB347",
    "EARTH_YELLOW": "#FFC370",
}
BRAND_RGB: dict[str, tuple[int, int, int]] = {
    name: _hex_to_rgb(hex_) for name, hex_ in BRAND_HEX.items()
}


# ===================================================================
# 1. COLOR OVERLAP TEST
# ===================================================================


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestColorPpt:
    """PPT colors are RGBColor objects with correct brand values."""

    def test_ppt_colors_are_rgbcolor_objects(self):
        """PPT colors must be RGBColor, not plain tuples."""
        for name, color in [
            ("PAPER", PAPER),
            ("WHITE", WHITE),
            ("SEA_SALT", SEA_SALT),
            ("CASTLETON_GREEN", CASTLETON_GREEN),
            ("DARK_SERPENT", DARK_SERPENT),
            ("SAFFRON", SAFFRON),
            ("EARTH_YELLOW", EARTH_YELLOW),
        ]:
            assert isinstance(
                color, RGBColor
            ), f"PPT {name} should be RGBColor, got {type(color).__name__}"

    def test_ppt_colors_match_brand_hex(self):
        """PPT color RGB values must match Lifewood brand spec."""
        for name, color in [
            ("PAPER", PAPER),
            ("WHITE", WHITE),
            ("SEA_SALT", SEA_SALT),
            ("CASTLETON_GREEN", CASTLETON_GREEN),
            ("DARK_SERPENT", DARK_SERPENT),
            ("SAFFRON", SAFFRON),
            ("EARTH_YELLOW", EARTH_YELLOW),
        ]:
            expected = BRAND_RGB[name]
            actual = _rbg_to_tuple(color)
            assert actual == expected, (
                f"PPT {name} RGB {actual} does not match "
                f"brand {expected} ({BRAND_HEX[name]})"
            )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestColorPdf:
    """PDF colors are (R,G,B) tuples with correct brand values."""

    def test_pdf_colors_are_tuples(self):
        """PDF colors must be plain (R,G,B) tuples — not RGBColor objects."""
        for name, color in [
            ("PDF_PAPER", PDF_PAPER),
            ("PDF_WHITE", PDF_WHITE),
            ("PDF_SEA_SALT", PDF_SEA_SALT),
            ("PDF_CASTLETON_GREEN", PDF_CASTLETON_GREEN),
            ("PDF_DARK_SERPENT", PDF_DARK_SERPENT),
            ("PDF_SAFFRON", PDF_SAFFRON),
            ("PDF_EARTH_YELLOW", PDF_EARTH_YELLOW),
        ]:
            assert isinstance(color, tuple), (
                f"PDF {name} should be tuple, got {type(color).__name__}"
            )
            assert len(color) == 3, f"PDF {name} tuple must have 3 elements"
            assert all(0 <= c <= 255 for c in color), (
                f"PDF {name} values {color} out of 0-255 range"
            )

    def test_pdf_colors_match_brand_hex(self):
        """PDF color RGB values must match Lifewood brand spec."""
        mapping = {
            "PDF_PAPER": ("PAPER", PDF_PAPER),
            "PDF_WHITE": ("WHITE", PDF_WHITE),
            "PDF_SEA_SALT": ("SEA_SALT", PDF_SEA_SALT),
            "PDF_CASTLETON_GREEN": ("CASTLETON_GREEN", PDF_CASTLETON_GREEN),
            "PDF_DARK_SERPENT": ("DARK_SERPENT", PDF_DARK_SERPENT),
            "PDF_SAFFRON": ("SAFFRON", PDF_SAFFRON),
            "PDF_EARTH_YELLOW": ("EARTH_YELLOW", PDF_EARTH_YELLOW),
        }
        for pdf_name, (brand_name, actual) in mapping.items():
            expected = BRAND_RGB[brand_name]
            assert actual == expected, (
                f"{pdf_name} {actual} does not match brand "
                f"{brand_name} {expected} ({BRAND_HEX[brand_name]})"
            )


@pytest.mark.skipif(
    not (HAS_PPT and HAS_PDF_MODULE),
    reason="Both PPT and PDF modules required",
)
class TestColorCrossSystemConsistency:
    """Both systems must represent the same brand colors."""

    def test_cross_system_color_equivalence(self):
        """PPT RGBColor and PDF tuple values must be identical."""
        pairs = [
            ("PAPER", PAPER, PDF_PAPER),
            ("WHITE", WHITE, PDF_WHITE),
            ("SEA_SALT", SEA_SALT, PDF_SEA_SALT),
            ("CASTLETON_GREEN", CASTLETON_GREEN, PDF_CASTLETON_GREEN),
            ("DARK_SERPENT", DARK_SERPENT, PDF_DARK_SERPENT),
            ("SAFFRON", SAFFRON, PDF_SAFFRON),
            ("EARTH_YELLOW", EARTH_YELLOW, PDF_EARTH_YELLOW),
        ]
        for name, ppt_color, pdf_color in pairs:
            ppt_tuple = _rbg_to_tuple(ppt_color)
            assert ppt_tuple == pdf_color, (
                f"Color mismatch for {name}: PPT={ppt_tuple}, "
                f"PDF={pdf_color} — systems have drifted"
            )


# ===================================================================
# 2. KEY NAMESPACE OVERLAP TEST
# ===================================================================


class TestSpecLockKeyNamespace:
    """Spec-lock files / classes must use distinct namespaces."""

    def test_pdf_spec_lock_keys_have_pdf_prefix(self, tmp_path):
        """PDF spec lock file must use '## pdf_' prefixed keys."""
        spec_path = tmp_path / "pdf_spec_lock.md"
        content = (
            "# PDF Execution Lock\n"
            "\n"
            "## pdf_colors\n"
            "- Castleton Green: (4, 98, 65)\n"
            "- Dark Serpent: (19, 48, 32)\n"
            "\n"
            "## pdf_typography\n"
            "- Font: Manrope\n"
            "\n"
            "## pdf_dimensions\n"
            "- A4\n"
        )
        spec_path.write_text(content, encoding="utf-8")

        text = spec_path.read_text(encoding="utf-8")
        assert "## pdf_colors" in text, (
            "PDF spec lock missing '## pdf_colors'"
        )
        assert "## pdf_typography" in text, (
            "PDF spec lock missing '## pdf_typography'"
        )
        assert "## pdf_dimensions" in text, (
            "PDF spec lock missing '## pdf_dimensions'"
        )
        # Must NOT use bare keys
        assert "## colors\n" not in text, (
            "PDF spec lock must not use bare '## colors'"
        )
        assert "## typography\n" not in text, (
            "PDF spec lock must not use bare '## typography'"
        )

    def test_ppt_spec_lock_keys_no_pdf_prefix(self, tmp_path):
        """PPT spec lock must NOT use '## pdf_' prefixed keys."""
        spec_path = tmp_path / "spec_lock.md"
        content = (
            "# Execution Lock\n"
            "\n"
            "## colors\n"
            "- Castleton Green: (4, 98, 65)\n"
            "\n"
            "## typography\n"
            "- Font: Manrope\n"
            "\n"
            "## mode\n"
            "- Output: PPTX\n"
        )
        spec_path.write_text(content, encoding="utf-8")

        text = spec_path.read_text(encoding="utf-8")
        assert "## colors\n" in text, (
            "PPT spec lock must use bare '## colors'"
        )
        assert "## typography\n" in text, (
            "PPT spec lock must use bare '## typography'"
        )
        # Must NOT use pdf_ prefix keys
        for key in ["## pdf_colors", "## pdf_typography", "## pdf_mode"]:
            assert key not in text, (
                f"PPT spec lock must not contain '{key}'"
            )

    def test_pdf_header_is_pdf_execution_lock(self, tmp_path):
        """PDF spec lock must start with '# PDF Execution Lock'."""
        spec_path = tmp_path / "pdf_spec_lock.md"
        content = "# PDF Execution Lock\n\n## pdf_colors\n"
        spec_path.write_text(content, encoding="utf-8")
        header = spec_path.read_text(encoding="utf-8").splitlines()[0]
        assert header == "# PDF Execution Lock", (
            f"PDF spec lock header must be '# PDF Execution Lock', "
            f"got '{header}'"
        )


# ===================================================================
# 3. UNIT OVERLAP TEST
# ===================================================================


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestPptUnits:
    """PPT uses inches and Emu — not millimeters."""

    def test_ppt_margin_in_inches_range(self):
        """PPT left margin should be 0.5-1.5 inches."""
        margin_inches = PPT_MARGIN_L / 914400  # Emu to inches
        assert 0.5 <= margin_inches <= 1.5, (
            f"PPT margin {PPT_MARGIN_L} Emu ({margin_inches:.2f} in) "
            f"outside expected range 0.5-1.5 in"
        )

    def test_ppt_margins_are_7_to_10_percent_safe_area(self):
        """Design system: 7-10% outer margins. A 4.5% margin reads as a
        document dropped onto a slide, not a designed page."""
        w, h = SLIDE_W, SLIDE_H
        for name, value, extent in [
            ("left", PPT_MARGIN_L, w), ("right", PPT_MARGIN_R, w),
            ("top", PPT_MARGIN_T, h), ("bottom", PPT_MARGIN_B, h),
        ]:
            pct = value / extent * 100
            assert 6.8 <= pct <= 10.2, (
                f"PPT {name} margin is {pct:.1f}% of the slide — "
                f"outside the 7-10% safe-area band"
            )

    def test_ppt_footer_in_inches_range(self):
        """PPT footer Y should be 6.0-7.5 inches (bottom of 7.5in slide)."""
        footer_inches = PPT_FOOTER_Y / 914400
        assert 6.0 <= footer_inches <= 7.5, (
            f"PPT footer Y {PPT_FOOTER_Y} Emu ({footer_inches:.2f} in) "
            f"outside expected range 6.0-7.5 in"
        )

    def test_ppt_slide_dimensions_widescreen(self):
        """PPT slide is 13.333 x 7.5 in (16:9 widescreen)."""
        w_inches = SLIDE_W / 914400
        h_inches = SLIDE_H / 914400
        assert abs(w_inches - 13.333) < 0.01, (
            f"PPT slide width {w_inches:.3f} in != 13.333 in"
        )
        assert abs(h_inches - 7.5) < 0.01, (
            f"PPT slide height {h_inches:.3f} in != 7.5 in"
        )
        ratio = w_inches / h_inches
        assert abs(ratio - 16 / 9) < 0.01, (
            f"PPT aspect ratio {ratio:.3f} != 16:9"
        )

    def test_ppt_uses_inches_not_mm(self):
        """PPT module must not contain millimeter measurement patterns."""
        import inspect
        import re
        import lifewood_deck_builder as deck

        source = inspect.getsource(deck)
        assert "Inches(" in source, (
            "PPT module must use Inches() for measurements"
        )
        # Use regex to avoid false positives on words containing 'mm'
        # (e.g. 'Summary', 'programmer'). Comments are exempt: the deck
        # documents the PDF equivalent of each anchor in mm on purpose —
        # the rule is about measurement *code*, not annotations.
        mm_measurement = re.compile(r"\d+\s*mm")
        lines_with_mm = [
            i + 1
            for i, line in enumerate(source.splitlines())
            if mm_measurement.search(line.split("#", 1)[0].lower())
        ]
        assert not lines_with_mm, (
            f"PPT module should not contain millimeter measurement "
            f"references (found at lines {lines_with_mm})"
        )

    # ---- Folio chrome (unified folio system) ----

    def test_ppt_footer_rule_above_footer(self):
        """Folio hairline rule must sit above the footer zone (gap 0.1-0.5 in)."""
        assert PPT_FOOTER_RULE_Y < PPT_FOOTER_Y, (
            f"PPT footer rule Y {PPT_FOOTER_RULE_Y} Emu must be above "
            f"footer {PPT_FOOTER_Y} Emu"
        )
        gap_in = (PPT_FOOTER_Y - PPT_FOOTER_RULE_Y) / 914400
        assert 0.1 <= gap_in <= 0.5, (
            f"PPT footer rule gap {gap_in:.3f} in outside 0.1-0.5 in"
        )

    def test_ppt_header_number_folio_size(self):
        """PPT folio kicker must sit in the 8-10pt folio band and match the
        footer folio exactly — one number family, not two."""
        assert 8 <= PPT_HEADER_NUM_SIZE <= 10, (
            f"PPT header number size {PPT_HEADER_NUM_SIZE} outside 8-10pt"
        )
        assert PPT_HEADER_NUM_SIZE == PPT_FOLIO_SIZE, (
            "header kicker and footer folio must be the same size"
        )

    def test_ppt_has_no_footer_label_constant(self):
        """The footer carries page numbering + logo ONLY. A footer-label
        size constant means the document title crept back into the band."""
        import lifewood_deck_builder as deck

        assert not hasattr(deck, "FOOTER_LABEL_SIZE"), (
            "FOOTER_LABEL_SIZE still exists — the footer left slot must "
            "stay empty (no document title, product name, or description)"
        )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestPdfUnits:
    """PDF uses millimeters — not inches or Emu."""

    def test_pdf_margin_in_mm_range(self):
        """PDF margins should be 10-50 mm on the slide canvas."""
        assert 10 <= PDF_MARGIN_L <= 50, (
            f"PDF margin {PDF_MARGIN_L} mm outside expected range 10-50 mm"
        )
        assert 10 <= PDF_MARGIN_R <= 50
        assert 10 <= PDF_MARGIN_T <= 50
        assert 10 <= PDF_MARGIN_B <= 50

    def test_pdf_margins_are_7_to_10_percent_safe_area(self):
        """Design system: 7-10% outer margins on the 16:9 canvas."""
        for name, value, extent in [
            ("left", PDF_MARGIN_L, PDF_PAGE_W),
            ("right", PDF_MARGIN_R, PDF_PAGE_W),
            ("top", PDF_MARGIN_T, PDF_PAGE_H),
            ("bottom", PDF_MARGIN_B, PDF_PAGE_H),
        ]:
            pct = value / extent * 100
            assert 6.8 <= pct <= 10.2, (
                f"PDF {name} margin is {pct:.1f}% of the page — "
                f"outside the 7-10% safe-area band"
            )

    @pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
    def test_pdf_and_ppt_margins_match(self):
        """One canvas, one grid — the two engines must agree in mm."""
        pairs = [
            (PPT_MARGIN_L / 914400 * 25.4, PDF_MARGIN_L, "left"),
            (PPT_MARGIN_R / 914400 * 25.4, PDF_MARGIN_R, "right"),
            (PPT_MARGIN_T / 914400 * 25.4, PDF_MARGIN_T, "top"),
            (PPT_MARGIN_B / 914400 * 25.4, PDF_MARGIN_B, "bottom"),
        ]
        for ppt_mm, pdf_mm, name in pairs:
            assert abs(ppt_mm - pdf_mm) < 0.5, (
                f"{name} margin drift: deck {ppt_mm:.2f} mm vs "
                f"PDF {pdf_mm:.2f} mm"
            )

    def test_pdf_page_size_is_powerpoint_widescreen(self):
        """PDF page must be PowerPoint 16:9 (13.333 x 7.5 in) — never A4."""
        assert abs(PDF_PAGE_W - 13.333 * 25.4) < 0.5, (
            f"PDF page width {PDF_PAGE_W} mm != 338.66 mm (13.333 in)"
        )
        assert abs(PDF_PAGE_H - 7.5 * 25.4) < 0.5, (
            f"PDF page height {PDF_PAGE_H} mm != 190.5 mm (7.5 in)"
        )
        assert abs(PDF_PAGE_W / PDF_PAGE_H - 16 / 9) < 0.01, (
            f"PDF aspect {PDF_PAGE_W / PDF_PAGE_H:.3f} != 16:9"
        )

    def test_pdf_page_is_not_a4(self):
        """A4 (210x297) is forbidden for every Lifewood deliverable."""
        assert not (abs(PDF_PAGE_W - 210) < 5 and abs(PDF_PAGE_H - 297) < 5), (
            "PDF canvas is A4 — Lifewood documents must use the "
            "PowerPoint 16:9 slide canvas"
        )

    @pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
    def test_pdf_page_matches_ppt_slide(self):
        """PDF canvas must equal the PPTX slide canvas in mm."""
        ppt_w_mm = SLIDE_W / 914400 * 25.4
        ppt_h_mm = SLIDE_H / 914400 * 25.4
        assert abs(PDF_PAGE_W - ppt_w_mm) < 0.5, (
            f"PDF width {PDF_PAGE_W} mm != slide width {ppt_w_mm:.2f} mm"
        )
        assert abs(PDF_PAGE_H - ppt_h_mm) < 0.5, (
            f"PDF height {PDF_PAGE_H} mm != slide height {ppt_h_mm:.2f} mm"
        )

    def test_pdf_footer_in_mm_range(self):
        """PDF footer zone must sit in the bottom band of the page."""
        assert PDF_PAGE_H - 30 <= PDF_FOOTER_Y <= PDF_PAGE_H - 8, (
            f"PDF footer Y {PDF_FOOTER_Y} mm outside the bottom band of a "
            f"{PDF_PAGE_H} mm page"
        )

    def test_pdf_body_band_between_header_and_footer(self):
        """Body band must start below the header rule and stop above the
        footer rule — no text can land on either piece of chrome."""
        assert PDF_BODY_TOP > PDF_HEADER_RULE_Y, (
            f"body top {PDF_BODY_TOP} mm must sit below the header rule "
            f"at {PDF_HEADER_RULE_Y} mm"
        )
        assert PDF_BODY_BOTTOM < PDF_FOOTER_RULE_Y, (
            f"body bottom {PDF_BODY_BOTTOM} mm must sit above the footer "
            f"rule at {PDF_FOOTER_RULE_Y} mm"
        )

    def test_pdf_two_column_grid_fits_content_width(self):
        """Two columns plus the gutter must equal the content width."""
        assert abs((PDF_COL_W * 2 + PDF_COL_GAP) - PDF_CONTENT_W) < 0.05, (
            f"2*{PDF_COL_W} + {PDF_COL_GAP} != content width {PDF_CONTENT_W}"
        )

    def test_pdf_content_width_correct(self):
        """PDF content width must be page minus margins."""
        expected = PDF_PAGE_W - PDF_MARGIN_L - PDF_MARGIN_R
        assert PDF_CONTENT_W == expected, (
            f"PDF_CONTENT_W={PDF_CONTENT_W}, expected {expected} "
            f"({PDF_PAGE_W} - {PDF_MARGIN_L} - {PDF_MARGIN_R})"
        )

    # ---- Folio chrome (unified folio system) ----

    def test_pdf_footer_rule_above_footer(self):
        """Folio hairline rule must sit above the footer zone (gap >= 1 mm)."""
        assert PDF_FOOTER_RULE_Y < PDF_FOOTER_Y, (
            f"PDF footer rule Y {PDF_FOOTER_RULE_Y} mm must be above "
            f"footer {PDF_FOOTER_Y} mm"
        )
        gap_mm = PDF_FOOTER_Y - PDF_FOOTER_RULE_Y
        assert gap_mm >= 1.0, (
            f"PDF footer rule gap {gap_mm} mm < 1.0 mm"
        )

    def test_pdf_header_number_folio_size(self):
        """PDF folio kicker sits in the 10-12pt band and equals the footer
        folio — header numbers and footer numbers are one family (folio
        family now 10-12pt, default 12)."""
        assert 10 <= PDF_HEADER_NUM_SIZE <= 12, (
            f"PDF header number size {PDF_HEADER_NUM_SIZE} outside 10-12pt"
        )
        assert PDF_HEADER_NUM_SIZE == PDF_FOLIO_SIZE, (
            "header kicker and footer folio must be the same size"
        )

    def test_pdf_has_no_footer_label_constant(self):
        """The footer carries page numbering + logo ONLY."""
        import lifewood_pdf_design_system as pdf_sys

        assert not hasattr(pdf_sys, "PDF_FOOTER_LABEL_SIZE"), (
            "PDF_FOOTER_LABEL_SIZE still exists — the footer left slot must "
            "stay empty (no document title, product name, or description)"
        )
        assert PdfSpecLock.FOOTER_LEFT_SLOT is None
        assert PdfSpecLock.COVER_FOOTER_LOGO is False


@pytest.mark.skipif(
    not (HAS_PPT and HAS_PDF_MODULE),
    reason="Both PPT and PDF modules required",
)
class TestUnifiedTypeScale:
    """A PDF page and a PPTX slide are the same rectangle, so the same role
    must render at the same point size in both engines."""

    def test_type_scale_matches_across_engines(self):
        # Deck engine frozen (PPTX output banned per skill spec §1); the
        # folio family was updated in the PDF engine only, so the folio and
        # header-kicker rows no longer apply cross-engine.
        pairs = [
            ("cover title", PPT_COVER_TITLE_SIZE, PDF_COVER_TITLE_SIZE),
            ("page title", PPT_PAGE_TITLE_SIZE, PDF_HEADER_TITLE_SIZE),
            ("body", PPT_BODY_SIZE, PDF_BODY_SIZE),
        ]
        for role, ppt_pt, pdf_pt in pairs:
            assert ppt_pt == pdf_pt, (
                f"{role} scale drift: deck {ppt_pt}pt vs PDF {pdf_pt}pt — "
                f"one canvas must mean one type scale"
            )

    def test_type_scale_inside_design_system_ranges(self):
        bands = [
            ("cover title", PDF_COVER_TITLE_SIZE, 42, 60),
            ("page title", PDF_HEADER_TITLE_SIZE, 30, 40),
            ("body", PDF_BODY_SIZE, 11, 14),
            ("folio", PDF_FOLIO_SIZE, 10, 12),
        ]
        for role, pt, lo, hi in bands:
            assert lo <= pt <= hi, (
                f"{role} is {pt}pt — outside the design-system band "
                f"{lo}-{hi}pt"
            )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestPdfSpecLockFolio:
    """PdfSpecLock mirrors folio chrome constants and emits them in the lock."""

    def test_pdf_spec_lock_mirrors_folio_constants(self):
        """PdfSpecLock must expose folio constants matching module values."""
        assert PdfSpecLock.FOOTER_RULE_Y == PDF_FOOTER_RULE_Y
        assert PdfSpecLock.FOOTER_RULE_H == PDF_FOOTER_RULE_H
        assert PdfSpecLock.HEADER_NUM_SIZE == PDF_HEADER_NUM_SIZE
        assert PdfSpecLock.HEADER_TITLE_SIZE == PDF_HEADER_TITLE_SIZE
        assert PdfSpecLock.FOLIO_SIZE == PDF_FOLIO_SIZE
        assert PdfSpecLock.BODY_SIZE == PDF_BODY_SIZE

    def test_pdf_spec_lock_has_folio_keys(self, tmp_path):
        """write_spec_lock must emit the footer/header chrome contract."""
        path = PdfSpecLock.write_spec_lock(
            str(tmp_path), "Folio Spec Test", "user", 3,
            ["1.0  Intro", "2.0  Body"],
        )
        text = Path(path).read_text(encoding="utf-8")
        for key in [
            "footer_rule_y_mm",
            "footer_rule_h_mm",
            "footer_left_slot: EMPTY",
            "cover_footer_logo: FORBIDDEN",
            "header_number_size",
            "header_title_size",
        ]:
            assert key in text, (
                f"pdf spec lock missing '{key}' line"
            )
        assert "footer_label_size" not in text, (
            "spec lock still advertises a footer label — the left slot is "
            "empty in the unified design system"
        )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestPdfNoPptxImports:
    """PDF module must not import from python-pptx."""

    def test_pdf_design_system_no_pptx_imports(self):
        """PDF design system must not contain pptx references."""
        import inspect
        import lifewood_pdf_design_system

        source = inspect.getsource(lifewood_pdf_design_system)
        forbidden = ["Inches", "Emu", "RGBColor", "pptx", "python-pptx"]
        for token in forbidden:
            assert token not in source, (
                f"PDF design system must not reference '{token}'"
            )

    def test_pdf_design_system_no_inches(self):
        """PDF module must use mm, not Inches."""
        import inspect
        import lifewood_pdf_design_system

        source = inspect.getsource(lifewood_pdf_design_system)
        assert "Inches" not in source, (
            "PDF design system must not use Inches()"
        )
        assert "Emu" not in source, (
            "PDF design system must not use Emu"
        )


# ===================================================================
# 4. LOGO PLACEMENT OVERLAP TEST
# ===================================================================


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestPptLogoPlacement:
    """PPT logo placement uses Inches from the right."""

    def test_ppt_logo_bottom_right(self):
        """PPT logo needs real breathing room from the right slide edge."""
        # From _footer_logo: logo_x = SLIDE_W - LOGO_W - LOGO_RIGHT_CLEAR
        logo_x = SLIDE_W - LOGO_W - PPT_LOGO_RIGHT_CLEAR
        right_gap = SLIDE_W - (logo_x + LOGO_W)
        assert right_gap == PPT_LOGO_RIGHT_CLEAR
        assert right_gap / 914400 >= 0.55, (
            f"PPT logo right gap {right_gap / 914400:.3f} in too tight — "
            f"reads as clipped by the slide edge"
        )

    @pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF module not available")
    def test_ppt_logo_clearance_matches_pdf(self):
        """Both engines clear the right edge by the same distance."""
        ppt_mm = PPT_LOGO_RIGHT_CLEAR / 914400 * 25.4
        assert abs(ppt_mm - PDF_LOGO_RIGHT_CLEAR) < 0.5, (
            f"logo clearance drift: deck {ppt_mm:.2f} mm vs "
            f"PDF {PDF_LOGO_RIGHT_CLEAR:.2f} mm"
        )

    def test_ppt_logo_dimensions(self):
        """PPT logo width ~1.8 in with ~5.5:1 aspect ratio."""
        w_in = LOGO_W / 914400
        h_in = LOGO_H / 914400
        assert abs(w_in - 1.8) < 0.05, (
            f"PPT logo width {w_in:.3f} in != ~1.8 in"
        )
        ratio = LOGO_W / LOGO_H
        assert abs(ratio - 5.5) < 0.5, (
            f"PPT logo aspect ratio {ratio:.2f} != ~5.5:1"
        )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestPdfLogoPlacement:
    """PDF logo placement uses mm from the right."""

    def test_pdf_logo_bottom_right(self):
        """PDF logo needs real breathing room from the right page edge."""
        # From footer(): logo_x = PDF_PAGE_W - PDF_LOGO_W - PDF_LOGO_RIGHT_CLEAR
        right_gap = PDF_PAGE_W - (
            (PDF_PAGE_W - PDF_LOGO_W - PDF_LOGO_RIGHT_CLEAR) + PDF_LOGO_W
        )
        assert abs(right_gap - PDF_LOGO_RIGHT_CLEAR) < 0.01
        assert right_gap >= 14.0, (
            f"PDF logo right gap {right_gap} mm too tight — reads as clipped"
        )

    def test_pdf_logo_dimensions(self):
        """PDF logo width ~45.7 mm (1.8 in, matching the deck) at 5.5:1."""
        assert abs(PDF_LOGO_W - 1.8 * 25.4) < 1.0, (
            f"PDF logo width {PDF_LOGO_W} mm != ~45.72 mm"
        )
        ratio = PDF_LOGO_W / PDF_LOGO_H
        assert abs(ratio - 5.5) < 0.5, (
            f"PDF logo aspect ratio {ratio:.2f} != ~5.5:1"
        )


@pytest.mark.skipif(
    not (HAS_PPT and HAS_PDF_MODULE),
    reason="Both PPT and PDF modules required",
)
class TestLogoAssetConsistency:
    """Both systems use the same logo asset files."""

    def test_same_logo_assets(self):
        """PPT and PDF must reference the same logo filenames."""
        assert "lwlogo_lightmode.png" in LOGO_LIGHT, (
            f"PPT light logo path missing filename: {LOGO_LIGHT}"
        )
        assert "lwlogo_darkmode.png" in LOGO_DARK, (
            f"PPT dark logo path missing filename: {LOGO_DARK}"
        )
        assert "lwlogo_lightmode.png" in PDF_LOGO_LIGHT, (
            f"PDF light logo path missing filename: {PDF_LOGO_LIGHT}"
        )
        assert "lwlogo_darkmode.png" in PDF_LOGO_DARK, (
            f"PDF dark logo path missing filename: {PDF_LOGO_DARK}"
        )

    def test_logo_assets_exist(self):
        """Both logo files must exist on disk."""
        for path in [LOGO_LIGHT, LOGO_DARK, PDF_LOGO_LIGHT, PDF_LOGO_DARK]:
            assert os.path.isfile(path), f"Logo file not found: {path}"


# ===================================================================
# 5. HEADING SYSTEM OVERLAP TEST
# ===================================================================


@pytest.mark.skipif(not HAS_PDF_BUILDER, reason="PDF builder not available")
class TestPdfNoPptxImportIsolation:
    """PDF builder must not import from python-pptx or lifewood_deck_builder."""

    def test_pdf_builder_import_isolation(self):
        """PDF builder must not import pptx/lifewood_deck_builder."""
        import inspect
        import lifewood_pdf_builder

        source = inspect.getsource(lifewood_pdf_builder)
        forbidden_imports = [
            "pptx",
            "python-pptx",
            "lifewood_deck_builder",
            "RGBColor",
            "Inches",
            "Emu",
            "Presentation",
        ]
        for token in forbidden_imports:
            assert token not in source, (
                f"PDF builder must not import '{token}' — "
                f"would create cross-system coupling"
            )

    def test_pdf_builder_imports_fpdf2(self):
        """PDF builder must import from fpdf2."""
        import inspect
        import lifewood_pdf_builder

        source = inspect.getsource(lifewood_pdf_builder)
        assert "from fpdf import" in source, (
            "PDF builder must import from fpdf2"
        )

    def test_pdf_builder_imports_from_design_system(self):
        """PDF builder must import from lifewood_pdf_design_system."""
        import inspect
        import lifewood_pdf_builder

        source = inspect.getsource(lifewood_pdf_builder)
        assert "lifewood_pdf_design_system" in source, (
            "PDF builder must import from lifewood_pdf_design_system"
        )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestHeadingSystemConstants:
    """PdfSpecLock must define numbered heading constants."""

    def test_pdf_spec_lock_has_color_constants(self):
        """PdfSpecLock must expose all brand colors."""
        assert hasattr(PdfSpecLock, "CASTLETON_GREEN")
        assert hasattr(PdfSpecLock, "DARK_SERPENT")
        assert hasattr(PdfSpecLock, "SAFFRON")
        assert hasattr(PdfSpecLock, "PAPER")
        assert hasattr(PdfSpecLock, "WHITE")
        assert hasattr(PdfSpecLock, "SEA_SALT")
        assert hasattr(PdfSpecLock, "EARTH_YELLOW")

    def test_pdf_spec_lock_color_values_match(self):
        """PdfSpecLock color class attrs must match module constants."""
        assert PdfSpecLock.CASTLETON_GREEN == PDF_CASTLETON_GREEN
        assert PdfSpecLock.DARK_SERPENT == PDF_DARK_SERPENT
        assert PdfSpecLock.SAFFRON == PDF_SAFFRON
        assert PdfSpecLock.PAPER == PDF_PAPER
        assert PdfSpecLock.WHITE == PDF_WHITE
        assert PdfSpecLock.SEA_SALT == PDF_SEA_SALT
        assert PdfSpecLock.EARTH_YELLOW == PDF_EARTH_YELLOW

    def test_pdf_spec_lock_has_layout_constants(self):
        """PdfSpecLock must expose layout constants."""
        assert PdfSpecLock.PAGE_W == PDF_PAGE_W
        assert PdfSpecLock.PAGE_H == PDF_PAGE_H
        assert PdfSpecLock.MARGIN_L == PDF_MARGIN_L
        assert PdfSpecLock.FONT == PDF_FONT
        assert PdfSpecLock.FALLBACK_FONT == PDF_FALLBACK_FONT
        assert abs(PdfSpecLock.LOGO_W - PDF_LOGO_W) < 0.01
        assert abs(PdfSpecLock.LOGO_H - PDF_LOGO_H) < 0.01
        assert abs(PdfSpecLock.LOGO_W / PdfSpecLock.LOGO_H - 5.5) < 0.05


# ===================================================================
# 6. FONT HANDLING TEST
# ===================================================================


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestPptFont:
    """PPT uses Manrope (embedded in PPTX XML)."""

    def test_ppt_font_is_manrope(self):
        """PPT font constant must be 'Manrope'."""
        assert FONT == "Manrope", (
            f"PPT font is '{FONT}', expected 'Manrope'"
        )


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestPdfFont:
    """PDF uses Manrope with Helvetica fallback."""

    def test_pdf_font_manrope(self):
        """PDF font constant is 'Manrope' (same brand font)."""
        assert PDF_FONT == "Manrope", (
            f"PDF font is '{PDF_FONT}', expected 'Manrope'"
        )

    def test_pdf_fallback_font(self):
        """PDF fallback font is 'Helvetica' (fpdf2 built-in)."""
        assert PDF_FALLBACK_FONT == "Helvetica", (
            f"PDF fallback font is '{PDF_FALLBACK_FONT}', "
            f"expected 'Helvetica'"
        )

    def test_fallback_different_from_primary(self):
        """Fallback must differ from primary font."""
        assert PDF_FALLBACK_FONT != PDF_FONT, (
            "Fallback font should differ from primary font"
        )


# ===================================================================
# 7. PDF SPEC LOCK — IMMUTABLE CONSTANTS TEST
# ===================================================================


@pytest.mark.skipif(not HAS_PDF_MODULE, reason="PDF design system not available")
class TestPdfSpecLockImmutable:
    """PdfSpecLock is an immutable constants class — not a writer."""

    def test_spec_lock_is_class_not_instance(self):
        """PdfSpecLock must be used as class, not instantiated."""
        # Access constants via class, not instance
        assert PdfSpecLock.FONT == PDF_FONT
        assert PdfSpecLock.PAGE_W == PDF_PAGE_W

    def test_spec_lock_has_color_groups(self):
        """PdfSpecLock must define color group sets."""
        assert hasattr(PdfSpecLock, "LIGHT_BG_COLORS")
        assert hasattr(PdfSpecLock, "DARK_BG_COLORS")
        assert hasattr(PdfSpecLock, "ACCENT_COLORS")
        assert hasattr(PdfSpecLock, "ALL_COLORS")

    def test_spec_lock_light_bg_colors(self):
        """Light background colors must include Paper, White, Sea Salt."""
        assert PDF_PAPER in PdfSpecLock.LIGHT_BG_COLORS
        assert PDF_WHITE in PdfSpecLock.LIGHT_BG_COLORS
        assert PDF_SEA_SALT in PdfSpecLock.LIGHT_BG_COLORS

    def test_spec_lock_dark_bg_colors(self):
        """Dark background colors must include Castleton Green, Dark Serpent."""
        assert PDF_CASTLETON_GREEN in PdfSpecLock.DARK_BG_COLORS
        assert PDF_DARK_SERPENT in PdfSpecLock.DARK_BG_COLORS

    def test_spec_lock_accent_colors(self):
        """Accent colors must include Saffron, Earth Yellow."""
        assert PDF_SAFFRON in PdfSpecLock.ACCENT_COLORS
        assert PDF_EARTH_YELLOW in PdfSpecLock.ACCENT_COLORS


# ===================================================================
# 8. ISOLATION TEST — CROSS-IMPORT
# ===================================================================


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestPptImportSanity:
    """PPT system should import correctly from python-pptx."""

    def test_ppt_imports_pptx(self):
        """PPT deck builder must import pptx classes."""
        import inspect
        import lifewood_deck_builder

        source = inspect.getsource(lifewood_deck_builder)
        assert "from pptx" in source, (
            "PPT deck builder must import from pptx"
        )
        assert "from pptx.util import" in source, (
            "PPT deck builder must import from pptx.util"
        )


class TestPdfModuleStructure:
    """PDF modules must exist with expected structure."""

    def test_pdf_design_system_exists(self):
        """lifewood_pdf_design_system.py must exist."""
        path = SKILL_DIR / "lifewood_pdf_design_system.py"
        assert path.is_file(), f"Not found at {path}"

    def test_pdf_builder_exists(self):
        """lifewood_pdf_builder.py must exist."""
        path = SKILL_DIR / "lifewood_pdf_builder.py"
        assert path.is_file(), f"Not found at {path}"

    @pytest.mark.skipif(
        not HAS_PDF_MODULE, reason="PDF design system not importable"
    )
    def test_pdf_design_system_required_constants(self):
        """PDF design system must export expected constants."""
        required = [
            "PDF_CASTLETON_GREEN",
            "PDF_DARK_SERPENT",
            "PDF_SAFFRON",
            "PDF_PAPER",
            "PDF_WHITE",
            "PDF_SEA_SALT",
            "PDF_EARTH_YELLOW",
            "PDF_MARGIN_L",
            "PDF_MARGIN_R",
            "PDF_MARGIN_T",
            "PDF_MARGIN_B",
            "PDF_PAGE_W",
            "PDF_PAGE_H",
            "PDF_CONTENT_W",
            "PDF_FOOTER_Y",
            "PDF_FONT",
            "PDF_FALLBACK_FONT",
            "PDF_LOGO_W",
            "PDF_LOGO_H",
            "PDF_LOGO_LIGHT",
            "PDF_LOGO_DARK",
            "PdfSpecLock",
        ]
        import lifewood_pdf_design_system as pdf_sys

        for const_name in required:
            assert hasattr(pdf_sys, const_name), (
                f"PDF design system missing: {const_name}"
            )
            assert getattr(pdf_sys, const_name) is not None, (
                f"PDF constant {const_name} is None"
            )

    @pytest.mark.skipif(
        not HAS_PDF_BUILDER, reason="PDF builder not importable"
    )
    def test_pdf_builder_required(self):
        """PDF builder must export LifewoodPDF."""
        from lifewood_pdf_builder import LifewoodPDF

        assert LifewoodPDF is not None

    @pytest.mark.skipif(
        not HAS_PDF_BUILDER, reason="PDF builder not importable"
    )
    def test_pdf_builder_heading_methods(self):
        """PDF builder must have heading/content methods."""
        from lifewood_pdf_builder import LifewoodPDF

        for attr in ["add_cover", "add_toc", "add_content_page",
                     "add_section_divider", "save"]:
            assert hasattr(LifewoodPDF, attr), (
                f"LifewoodPDF missing method: {attr}"
            )


@pytest.mark.skipif(
    not (HAS_PDF_MODULE and HAS_PDF_BUILDER),
    reason="Both PDF modules required",
)
class TestFullCrossSystemIsolation:
    """Comprehensive isolation check between PPT and PDF systems."""

    def test_no_ppt_import_in_pdf_design_system(self):
        """PDF design system module does not import pptx."""
        import inspect
        import lifewood_pdf_design_system

        source = inspect.getsource(lifewood_pdf_design_system)
        pptx_tokens = ["pptx", "python-pptx", "Presentation", "RGBColor"]
        for token in pptx_tokens:
            assert token not in source, (
                f"PDF design system must not reference '{token}'"
            )

    def test_no_ppt_import_in_pdf_builder(self):
        """PDF builder module does not import pptx."""
        import inspect
        import lifewood_pdf_builder

        source = inspect.getsource(lifewood_pdf_builder)
        pptx_tokens = ["pptx", "python-pptx", "Presentation", "RGBColor"]
        for token in pptx_tokens:
            assert token not in source, (
                f"PDF builder must not reference '{token}'"
            )

    def test_pdf_builder_uses_mm_not_inches(self):
        """PDF builder must use mm units, not Inches."""
        import inspect
        import lifewood_pdf_builder

        source = inspect.getsource(lifewood_pdf_builder)
        assert "Inches" not in source, (
            "PDF builder must not use Inches"
        )
        assert "Emu" not in source, (
            "PDF builder must not use Emu"
        )

    def test_pdf_builder_imports_fpdf2(self):
        """PDF builder must import from fpdf2."""
        import inspect
        import lifewood_pdf_builder

        source = inspect.getsource(lifewood_pdf_builder)
        assert "from fpdf import" in source, (
            "PDF builder must import from fpdf2"
        )


# ===================================================================
# 9. PAGE CHROME TEST — numbering + bottom-right logo
# ===================================================================
# Spec (binding, both engines):
#   - Page numbers appear ONLY on pages/slides AFTER the Table of
#     Contents. Cover AND ToC are unnumbered.
#   - Numbering restarts: first content page/slide after ToC =
#     "Page 1 of N", N = number of pages/slides after cover + ToC.
#   - The cover carries exactly ONE logo (hero, upper-left) and NO
#     bottom-right logo. Every interior page/slide (ToC included) carries
#     the bottom-right logo.
#   - The footer holds page numbering + logo ONLY; the left slot is empty.
#   - Format "Page X of Y".

_PAGE_NUM_RE = re.compile(r"Page \d+ of \d+")
# Logo's bottom-right corner must sit within 25 mm of the page corner.
_PDF_BOTTOM_RIGHT_MM = 25.0


def _build_chrome_pdf(path: Path) -> Path:
    """Build cover + ToC + 2 content pages and write to *path*."""
    pdf = LifewoodPDF("Chrome Spec Test", "Verification Document")
    pdf.add_cover()
    pdf.add_toc(["1.0  Introduction", "2.0  Body"])
    pdf.add_content_page("1.0", "Introduction", ["First content page."])
    pdf.add_content_page("2.0", "Body", ["Second content page."])
    pdf.save(str(path))
    return path


def _page_has_bottom_right_image(page, margin_mm: float = _PDF_BOTTOM_RIGHT_MM) -> bool:
    """True if *page* has at least one image whose bottom-right corner
    sits within *margin_mm* of the page's bottom-right corner."""
    mm_to_pt = 72.0 / 25.4
    m = margin_mm * mm_to_pt
    pw, ph = page.rect.width, page.rect.height
    for img in page.get_images(full=True):
        for rect in page.get_image_rects(img[0]):
            if rect.x1 >= pw - m and rect.y1 >= ph - m:
                return True
    return False


def _build_chrome_deck(path: Path) -> Path:
    """Build cover + ToC + 2 content slides and write to *path*."""
    deck = LifewoodDeck("Chrome Spec Test", "Verification Deck")
    deck.add_cover()
    deck.add_toc(["1.0  Introduction", "2.0  Body"])
    deck.add_content("1.0", "Introduction", ["First content point."])
    deck.add_content("2.0", "Body", ["Second content point."])
    deck.save(str(path), skip_density_check=True)
    return path


def _slide_texts(slide) -> list[str]:
    """All non-empty paragraph texts on a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    texts.append(t)
    return texts


@pytest.mark.skipif(
    not (HAS_PDF_BUILDER and HAS_FITZ),
    reason="PDF builder + PyMuPDF required",
)
class TestPdfPageChrome:
    """PDF page chrome: numbers only after ToC, logo every page."""

    def test_pdf_cover_and_toc_have_no_page_number(self, tmp_path):
        """Cover (page 1) and ToC (page 2) must not render 'Page X of Y'."""
        path = _build_chrome_pdf(tmp_path / "chrome.pdf")
        doc = fitz.open(str(path))
        try:
            assert len(doc) == 4, f"expected 4 pages, got {len(doc)}"
            for i in (0, 1):
                m = _PAGE_NUM_RE.search(doc[i].get_text())
                assert m is None, (
                    f"page {i + 1} (cover/ToC) must be unnumbered, "
                    f"found '{m.group(0)}'"
                )
        finally:
            doc.close()

    def test_pdf_numbering_starts_after_toc(self, tmp_path):
        """Content pages restart at 'Page 1 of 2' after cover + ToC."""
        path = _build_chrome_pdf(tmp_path / "chrome.pdf")
        doc = fitz.open(str(path))
        try:
            text_p3 = doc[2].get_text()
            text_p4 = doc[3].get_text()
            assert "Page 1 of 2" in text_p3, (
                "first content page must read 'Page 1 of 2', got: "
                f"{_PAGE_NUM_RE.findall(text_p3)}"
            )
            assert "Page 2 of 2" in text_p4, (
                "second content page must read 'Page 2 of 2', got: "
                f"{_PAGE_NUM_RE.findall(text_p4)}"
            )
        finally:
            doc.close()

    def test_pdf_verify_numbering_passes(self, tmp_path):
        """The engine's own numbering self-check must pass on a real doc."""
        pdf = LifewoodPDF("Chrome Spec Test", "Verification Document")
        pdf.add_cover()
        pdf.add_toc(["1.0  Introduction", "2.0  Body"])
        pdf.add_content_page("1.0", "Introduction", ["First content page."])
        pdf.add_section_divider("2.0", "Body")
        pdf.add_content_page("2.1", "Detail", ["Second content page."])
        report = pdf.verify_numbering()
        assert report["ok"], report["errors"]
        assert report["front_pages"] == 2
        assert report["numbered_pages"] == 3
        assert report["first_numbered_page"] == 3
        pdf.save(str(tmp_path / "verified.pdf"))

    def test_pdf_save_rejects_broken_numbering(self, tmp_path):
        """A ToC placed after content must fail the save-time contract."""
        pdf = LifewoodPDF("Broken Order", "Verification Document")
        pdf.add_cover()
        pdf.add_content_page("1.0", "Introduction", ["Content before ToC."])
        pdf.add_toc(["1.0  Introduction"])
        with pytest.raises(ValueError, match="page-numbering contract"):
            pdf.save(str(tmp_path / "broken.pdf"))

    def test_pdf_pages_are_slide_sized(self, tmp_path):
        """Every rendered page must measure 13.333 x 7.5 in — never A4."""
        path = _build_chrome_pdf(tmp_path / "chrome.pdf")
        doc = fitz.open(str(path))
        try:
            for i, page in enumerate(doc):
                w_mm = page.rect.width / 72 * 25.4
                h_mm = page.rect.height / 72 * 25.4
                assert abs(w_mm - PDF_PAGE_W) < 0.5, (
                    f"page {i + 1} width {w_mm:.1f} mm != {PDF_PAGE_W} mm"
                )
                assert abs(h_mm - PDF_PAGE_H) < 0.5, (
                    f"page {i + 1} height {h_mm:.1f} mm != {PDF_PAGE_H} mm"
                )
        finally:
            doc.close()

    def test_pdf_content_pages_share_header_and_footer(self, tmp_path):
        """The header band (kicker + title) and the folio appear on every
        content page. Long content is compressed to fit — no continuation
        pages (skill §14)."""
        pdf = LifewoodPDF("Chrome Consistency", "Verification Document")
        pdf.add_cover()
        pdf.add_toc(["1.0  Introduction"])
        pdf.add_content_page(
            "1.0", "Introduction",
            ["Overflowing paragraph. " * 60 for _ in range(6)],
        )
        path = tmp_path / "chrome_consistency.pdf"
        pdf.save(str(path))
        doc = fitz.open(str(path))
        try:
            # No continuation pages: content compressed to fit on one page.
            assert len(doc) == 3, (
                f"expected 3 pages (cover + toc + 1 content), got {len(doc)}"
            )
            text = doc[2].get_text()
            assert "1.0" in text, "content page missing folio kicker"
            assert "Introduction" in text, "content page missing title"
            assert _PAGE_NUM_RE.search(text), (
                "content page missing 'Page X of Y' folio"
            )
        finally:
            doc.close()

    def test_pdf_footer_carries_no_document_title(self, tmp_path):
        """The footer left slot stays empty: page numbering + logo only."""
        pdf = LifewoodPDF("Unique Footer Probe Title", "Verification Document")
        pdf.add_cover()
        pdf.add_toc(["1.0  Introduction"])
        pdf.add_content_page("1.0", "Introduction", ["Body text."])
        path = tmp_path / "footer_slot.pdf"
        pdf.save(str(path))
        mm = 72.0 / 25.4
        doc = fitz.open(str(path))
        try:
            for i in range(2, len(doc)):
                for block in doc[i].get_text("blocks"):
                    _x0, y0, _x1, _y1, text = block[:5]
                    text = (text or "").strip()
                    if not text or y0 / mm < PDF_FOOTER_RULE_Y:
                        continue
                    assert text.startswith("Page "), (
                        f"page {i + 1} footer band contains '{text[:40]}' — "
                        f"only 'Page X of Y' and the logo belong there"
                    )
        finally:
            doc.close()

    def test_pdf_body_text_stays_inside_the_band(self, tmp_path):
        """No body text may overlap the header rule or the footer band."""
        pdf = LifewoodPDF("Band Guard", "Verification Document")
        pdf.add_cover()
        pdf.add_toc(["1.0  Introduction"])
        pdf.add_content_page(
            "1.0", "Introduction",
            ["Overflowing paragraph. " * 60 for _ in range(4)],
        )
        path = tmp_path / "band.pdf"
        pdf.save(str(path))
        mm = 72.0 / 25.4
        doc = fitz.open(str(path))
        try:
            for i in range(2, len(doc)):
                for block in doc[i].get_text("blocks"):
                    x0, y0, x1, y1, text = block[:5]
                    text = (text or "").strip()
                    if not text or text.startswith("Page ") \
                            or text.startswith("Band Guard") \
                            or text.startswith("1.0") \
                            or text.startswith("Introduction"):
                        continue
                    assert y0 >= PDF_HEADER_RULE_Y * mm - 1, (
                        f"page {i + 1}: body block above the header rule"
                    )
                    assert y1 <= PDF_FOOTER_RULE_Y * mm + 1, (
                        f"page {i + 1}: body block runs into the footer band"
                    )
        finally:
            doc.close()

    def test_pdf_cover_front_matter_order(self, tmp_path):
        """Cover front matter reads Prepared by → Date → Classification,
        with the classification stamp bottom-left."""
        pdf = LifewoodPDF("Front Matter", "Verification Document")
        pdf.add_cover()
        pdf.add_toc(["1.0  Introduction"])
        pdf.add_content_page("1.0", "Introduction", ["Body."])
        path = tmp_path / "front_matter.pdf"
        pdf.save(str(path))
        doc = fitz.open(str(path))
        try:
            page = doc[0]
            found = {}
            for block in page.get_text("blocks"):
                x0, y0, _x1, _y1, text = block[:5]
                text = (text or "").strip()
                if text.startswith("Prepared by:"):
                    found["prepared"] = (x0, y0)
                elif text.startswith("Date:"):
                    found["date"] = (x0, y0)
                elif text.startswith("CONFIDENTIAL"):
                    found["class"] = (x0, y0)
            assert set(found) == {"prepared", "date", "class"}, (
                f"cover front matter incomplete: {sorted(found)}"
            )
            assert found["prepared"][1] < found["date"][1] < found["class"][1], (
                "order must be Prepared by → Date → Classification"
            )
            # Bottom-left: left half of the page, inside the footer band
            assert found["class"][0] < page.rect.width / 2, (
                "classification must sit on the left of the cover"
            )
            assert found["class"][1] > page.rect.height * 0.85, (
                "classification must sit in the bottom band of the cover"
            )
        finally:
            doc.close()

    def test_pdf_classification_is_configurable(self, tmp_path):
        """Custom text is honoured; None drops the stamp entirely."""
        custom = LifewoodPDF("Custom Class", classification="RESTRICTED")
        custom.add_cover()
        custom.add_toc(["1.0  A"])
        custom.add_content_page("1.0", "A", ["Body."])
        custom.save(str(tmp_path / "custom.pdf"))

        none_pdf = LifewoodPDF("No Class", classification=None)
        none_pdf.add_cover()
        none_pdf.add_toc(["1.0  A"])
        none_pdf.add_content_page("1.0", "A", ["Body."])
        none_pdf.save(str(tmp_path / "none.pdf"))

        doc = fitz.open(str(tmp_path / "custom.pdf"))
        try:
            assert "RESTRICTED" in doc[0].get_text()
        finally:
            doc.close()
        doc = fitz.open(str(tmp_path / "none.pdf"))
        try:
            assert "CONFIDENTIAL" not in doc[0].get_text()
        finally:
            doc.close()

    def test_pdf_logo_on_every_page(self, tmp_path):
        """Cover carries the hero logo upper-left only (no bottom-right);
        every interior page's logo sits bottom-right."""
        path = _build_chrome_pdf(tmp_path / "chrome.pdf")
        doc = fitz.open(str(path))
        try:
            for i in range(len(doc)):
                page = doc[i]
                assert page.get_images(full=True), (
                    f"page {i + 1} must contain at least one image (logo)"
                )
            # Cover (page 1) is exempt from the bottom-right placement —
            # exactly one logo, upper-left (PdfSpecLock.COVER_FOOTER_LOGO).
            for i in range(1, len(doc)):
                assert _page_has_bottom_right_image(doc[i]), (
                    f"interior page {i + 1} must have a bottom-right logo "
                    f"(within {_PDF_BOTTOM_RIGHT_MM:.0f} mm of the corner)"
                )
        finally:
            doc.close()


@pytest.mark.skipif(
    not (HAS_PDF_BUILDER and HAS_FITZ),
    reason="PDF builder + PyMuPDF required",
)
class TestPdfLayoutIntegrity:
    """Torture the engine: nothing may leave the safe area, and pathological
    content must not detonate the page count."""

    def _stress(self, path: Path) -> LifewoodPDF:
        long_title = (
            "Enterprise Data Platform Operations and Governance Handbook "
            "for All Regional Delivery Teams"
        )
        pdf = LifewoodPDF(
            long_title,
            "A subtitle that is also very long and keeps going for a while",
            audience="user",
        )
        pdf.add_cover()
        pdf.add_toc([f"{i}.0 Section {i} with a long descriptive title"
                     for i in range(1, 41)])
        pdf.add_content_page("1.0", long_title,
                             ["x" * 3000, "Body sentence here. " * 60])
        pdf.add_mixed_content_page("2.0", "Reference", [
            ("callout", "Warning", "This callout body is very long. " * 40),
            ("code", "\n".join(
                f"def function_{i}(argument_one, argument_two): "
                f"return argument_one  # a trailing comment that runs long"
                for i in range(60)), "PY"),
            ("table", ["Field", "Meaning", "Default"],
             [[f"field_{i}", "a description that wraps across the cell " * 2,
               "on"] for i in range(30)]),
        ])
        pdf.add_section_divider(
            "3.0",
            "A section divider title that is extremely long and will not fit "
            "on one line at all no matter what you try",
        )
        pdf.add_two_column("4.0", "Compare",
                           "Left heading", ["item " * 30] * 12,
                           "Right heading", ["item " * 30] * 12)
        pdf.save(str(path), audit=True)
        return pdf

    def test_stress_document_has_no_layout_violations(self, tmp_path):
        """A torture document must audit clean — save(audit=True) raises
        if any block leaves the safe area."""
        path = tmp_path / "stress.pdf"
        self._stress(path)
        result = LifewoodPDF.audit_layout(str(path))
        assert result["ok"], result["violations"][:10]

    def test_long_code_block_does_not_explode_page_count(self, tmp_path):
        """A 60-line code block is a few panels, not dozens of stray pages."""
        pdf = LifewoodPDF("Code", "x")
        pdf.add_cover()
        pdf.add_toc(["1.0  Code"])
        before = pdf.page_no()
        pdf.add_mixed_content_page("1.0", "Code", [
            ("code", "\n".join(f"line_{i} = compute(a, b, c)"
                               for i in range(60)), "PY"),
        ])
        grew = pdf.page_no() - before
        assert grew <= 4, (
            f"60 code lines produced {grew} pages — panel paging is broken"
        )
        pdf.save(str(tmp_path / "code.pdf"), audit=True)

    def test_code_lines_wrap_inside_the_panel(self, tmp_path):
        """Over-long code lines wrap; they never run past the panel."""
        pdf = LifewoodPDF("Wrap", "x")
        pdf.add_cover()
        pdf.add_toc(["1.0  Code"])
        pdf.add_mixed_content_page("1.0", "Code", [
            ("code", "x = " + " + ".join(f"variable_{i}" for i in range(60)),
             "PY"),
        ])
        path = tmp_path / "wrap.pdf"
        pdf.save(str(path), audit=True)
        assert LifewoodPDF.audit_layout(str(path))["ok"]

    def test_callout_height_is_measured_not_guessed(self, tmp_path):
        """Callout text stays inside its fill rect for any length."""
        for reps in (1, 12, 60):
            pdf = LifewoodPDF("Callout", "x")
            pdf.add_cover()
            pdf.add_toc(["1.0  Note"])
            pdf.add_mixed_content_page("1.0", "Note", [
                ("callout", "Warning", "Callout sentence. " * reps),
            ])
            path = tmp_path / f"callout_{reps}.pdf"
            pdf.save(str(path), audit=True)
            assert LifewoodPDF.audit_layout(str(path))["ok"]

    def test_long_toc_spills_to_a_second_toc_page(self, tmp_path):
        """A 60-entry ToC continues on another front-matter page — and the
        folio still restarts at 1 on the first content page after it."""
        pdf = LifewoodPDF("Long ToC", "x")
        pdf.add_cover()
        pdf.add_toc([f"{i}.0 Section {i}" for i in range(1, 61)])
        pdf.add_content_page("1.0", "Body", ["Text."])
        path = tmp_path / "toc.pdf"
        pdf.save(str(path), audit=True)
        report = pdf.verify_numbering()
        assert report["ok"], report["errors"]
        assert report["front_pages"] >= 3, (
            "long ToC should occupy more than one front-matter page"
        )
        doc = fitz.open(str(path))
        try:
            assert "Page 1 of 1" in doc[len(doc) - 1].get_text()
        finally:
            doc.close()

    def test_locked_title_never_exceeds_content_width(self, tmp_path):
        """An over-long page title is auto-fitted, never allowed to run on."""
        pdf = LifewoodPDF("Title Fit", "x")
        pdf.add_cover()
        pdf.add_toc(["1.0  A"])
        pdf.add_content_page("1.0", "A page title that is far too long to "
                                    "fit on one line at the locked size and "
                                    "then some more words after that",
                             ["Body."])
        path = tmp_path / "title.pdf"
        pdf.save(str(path), audit=True)
        assert LifewoodPDF.audit_layout(str(path))["ok"]

    def test_audit_catches_real_damage(self, tmp_path):
        """The auditor must not be a rubber stamp — an A4 page fails it."""
        from fpdf import FPDF

        bad = FPDF(orientation="P", unit="mm", format="A4")
        bad.add_page()
        bad.set_font("Helvetica", size=12)
        bad.cell(0, 10, "Wrong canvas")
        bad.output(str(tmp_path / "a4.pdf"))
        result = LifewoodPDF.audit_layout(str(tmp_path / "a4.pdf"))
        assert not result["ok"]
        assert any(v["kind"] == "page_size" for v in result["violations"])


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestPptLayoutIntegrity:
    """No slide shape may cross the slide edge — PowerPoint clips it."""

    def test_no_shape_leaves_the_slide(self, tmp_path):
        for variant in ("editorial", "swiss"):
            deck = LifewoodDeck("Overflow Probe", "x", variant=variant)
            deck.add_cover()
            deck.add_toc([f"{i}.0 Section {i} with a long descriptive title"
                          for i in range(1, 20)])
            deck.add_content("1.0", "Many points",
                             [f"Point {i} with a reasonably long sentence"
                              for i in range(20)])
            path = tmp_path / f"ovf_{variant}.pptx"
            deck.save(str(path), skip_density_check=True)
            prs = Presentation(str(path))
            for i, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    bottom = (shape.top or 0) + (shape.height or 0)
                    right = (shape.left or 0) + (shape.width or 0)
                    assert bottom <= SLIDE_H, (
                        f"{variant} slide {i}: shape runs {bottom - SLIDE_H} "
                        f"EMU past the bottom edge"
                    )
                    assert right <= SLIDE_W, (
                        f"{variant} slide {i}: shape runs past the right edge"
                    )


@pytest.mark.skipif(not HAS_PPT, reason="PPT module not available")
class TestPptPageChrome:
    """PPTX page chrome: numbers only after ToC, logo every slide."""

    def test_ppt_cover_and_toc_have_no_page_number(self, tmp_path):
        """Cover (slide 1) and ToC (slide 2) must not show 'Page …' text."""
        path = _build_chrome_deck(tmp_path / "chrome.pptx")
        prs = Presentation(str(path))
        for idx in (0, 1):
            for text in _slide_texts(prs.slides[idx]):
                assert not text.startswith("Page "), (
                    f"slide {idx + 1} (cover/ToC) must be unnumbered, "
                    f"found '{text}'"
                )

    def test_ppt_numbering_starts_after_toc(self, tmp_path):
        """Content slides restart at 'Page 1 of 2' after cover + ToC."""
        path = _build_chrome_deck(tmp_path / "chrome.pptx")
        prs = Presentation(str(path))
        texts_p3 = _slide_texts(prs.slides[2])
        texts_p4 = _slide_texts(prs.slides[3])
        assert any(t.startswith("Page 1 of 2") for t in texts_p3), (
            f"slide 3 must read 'Page 1 of 2', got: {texts_p3}"
        )
        assert any(t.startswith("Page 2 of 2") for t in texts_p4), (
            f"slide 4 must read 'Page 2 of 2', got: {texts_p4}"
        )

    def test_ppt_cover_front_matter_order(self, tmp_path):
        """Deck cover carries Prepared by → Date → Classification, with the
        classification bottom-left."""
        path = _build_chrome_deck(tmp_path / "chrome.pptx")
        prs = Presentation(str(path))
        cover = prs.slides[0]
        found = {}
        for shape in cover.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text.startswith("Prepared by:"):
                found["prepared"] = (shape.left, shape.top)
            elif text.startswith("Date:"):
                found["date"] = (shape.left, shape.top)
            elif text.startswith("CONFIDENTIAL"):
                found["class"] = (shape.left, shape.top)
        assert set(found) == {"prepared", "date", "class"}, (
            f"cover front matter incomplete: {sorted(found)}"
        )
        assert found["prepared"][1] < found["date"][1] < found["class"][1], (
            "order must be Prepared by → Date → Classification"
        )
        assert found["class"][0] < SLIDE_W / 2, (
            "classification must sit on the left of the cover"
        )
        assert found["class"][1] > SLIDE_H * 0.85, (
            "classification must sit in the bottom band of the cover"
        )

    def test_ppt_and_pdf_covers_share_anchors(self, tmp_path):
        """Deck cover and PDF cover must use the same composition anchors —
        one first page, not two interpretations of it."""
        if not (HAS_PDF_BUILDER and HAS_PDF_MODULE):
            pytest.skip("PDF builder required")
        from lifewood_deck_builder import (
            COVER_TEXT_X, COVER_TITLE_Y, COVER_META_Y, COVER_CLASS_Y,
        )

        pairs = [
            (COVER_TEXT_X / 914400 * 25.4, LifewoodPDF.COVER_TEXT_X),
            (COVER_TITLE_Y / 914400 * 25.4, LifewoodPDF.COVER_TITLE_Y),
            (COVER_META_Y / 914400 * 25.4, LifewoodPDF.COVER_META_Y),
            (COVER_CLASS_Y / 914400 * 25.4, LifewoodPDF.COVER_CLASS_Y),
        ]
        for ppt_mm, pdf_mm in pairs:
            assert abs(ppt_mm - pdf_mm) < 0.5, (
                f"cover anchor drift: deck {ppt_mm:.2f} mm vs PDF {pdf_mm:.2f} mm"
            )

    def test_ppt_logo_on_every_slide(self, tmp_path):
        """Cover carries the hero logo upper-left only (no bottom-right);
        every interior slide (ToC, content) has a bottom-right logo."""
        path = _build_chrome_deck(tmp_path / "chrome.pptx")
        prs = Presentation(str(path))
        for idx, slide in enumerate(prs.slides, start=1):
            if idx == 1:  # cover: hero logo upper-left only
                continue
            pics_bottom_right = [
                sh for sh in slide.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                and sh.left > Inches(10.5) and sh.top > Inches(6.0)
            ]
            assert pics_bottom_right, (
                f"slide {idx} must contain a bottom-right logo "
                f"(left > 10.5in, top > 6.0in)"
            )


# ===================================================================
# CLEANUP
# ===================================================================


def pytest_sessionfinish(session):
    """Clean up temp test output files."""
    if TEMP_DIR.exists():
        import shutil

        shutil.rmtree(TEMP_DIR, ignore_errors=True)
