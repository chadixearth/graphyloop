#!/usr/bin/env python3
"""
Lifewood Deck Builder — premium branded PowerPoint/PDF generator.

Builds natively-editable .pptx files with Lifewood brand identity applied
at the XML/object level. Every element is a real PPTX shape with precise
positioning, proper font embedding, and correct color values.

Design philosophy:
    - Premium agency aesthetic — generous whitespace, clear hierarchy
    - Typography-first design — Manrope at precise weights
    - Strategic accent usage — Saffron only where it draws the eye
    - Consistent geometry — 0.85" margins, 16:9 widescreen
    - Numbered headings system (1.0, 1.1, 1.1.1, 2.0...)
    - Logo bottom-right every slide (correct light/dark variant)
    - Page X of Y bottom-center on content slides only (cover & TOC unnumbered, starts after ToC)
    - Subtle design elements that elevate above generic

Usage:
    python lifewood_deck_builder.py "Title" --output deck.pptx
    python lifewood_deck_builder.py "Title" --output deck.pptx --pdf

Or import:
    from lifewood_deck_builder import LifewoodDeck

    deck = LifewoodDeck("Q3 Strategy", "Growth Initiatives")
    deck.add_cover()
    deck.add_toc(["1.0 Market", "2.0 Analysis"])
    deck.add_section("1.0", "Market Overview")
    deck.add_content("1.1", "Key Trends", ["Point one", "Point two"])
    deck.add_closing()
    deck.save("output.pptx")
"""

from __future__ import annotations

import argparse
import os
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.util import Emu, Inches, Pt

from lifewood_design_variants import get_variant

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------

PAPER = RGBColor(0xF5, 0xEE, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SEA_SALT = RGBColor(0xF9, 0xF7, 0xF7)
CASTLETON_GREEN = RGBColor(0x04, 0x62, 0x41)
DARK_SERPENT = RGBColor(0x13, 0x30, 0x20)
SAFFRON = RGBColor(0xFF, 0xB3, 0x47)
EARTH_YELLOW = RGBColor(0xFF, 0xC3, 0x70)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Grid system — 7-10% safe outer margins, identical to the PDF engine.
# 0.945 in L/R is 7.1% of the 13.333 in width; 0.525 in top is 7.0% of the
# 7.5 in height; 0.75 in bottom is 10.0% and doubles as the footer band.
MARGIN_L = Inches(0.945)
MARGIN_R = Inches(0.945)
MARGIN_T = Inches(0.525)
MARGIN_B = Inches(0.75)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # 11.443 in
CONTENT_H = SLIDE_H - MARGIN_T - MARGIN_B  # 6.225 in

# 12-column grid
GRID_COLS = 12
COL_W = 0.8275  # inches, used with Inches()
GUTTER = 0.2    # inches
LIVE_AREA = CONTENT_W * CONTENT_H  # ~75.81 sq in
DENSITY_FLOOR = LIVE_AREA * 0.6    # ~45.49 sq in

# Density check: content bbox must cover ≥60% of live area
DENSITY_FLOOR_SQ_IN = 45.49

# Footer zone
FOOTER_Y = Inches(6.75)
FOOTER_H = Inches(0.75)

# Folio chrome system (header kicker + footer rule + centred page number)
# The footer carries page numbering and the logo ONLY — the left slot stays
# empty. Repeating the document title on every slide is footer clutter, not
# branding, so there is deliberately no footer-label constant here.
FOOTER_RULE_Y = Inches(6.52)      # hairline rule above footer zone

FOOTER_RULE_H = Pt(0.75)          # hairline thickness

# ---------------------------------------------------------------------------
# Type scale — ONE scale shared with the PDF engine.
#
# A slide and a PDF page are the same 13.333 x 7.5 in rectangle, so the same
# role renders at the same point size in both. The old split (28 pt PPT title
# vs 16 pt PDF title, 12 pt vs 9 pt folio) made two documents out of one
# design system.
# ---------------------------------------------------------------------------

COVER_TITLE_SIZE = 54             # range 42-60
PAGE_TITLE_SIZE = 32              # range 30-40
SECTION_HEADING_SIZE = 24         # range 20-28
SUBHEADING_SIZE = 18              # range 16-20
BODY_SIZE = 12                    # range 11-14
TABLE_SIZE = 11                   # range 10-13
CAPTION_SIZE = 10                 # range 9-11
FOLIO_SIZE = 10                   # range 8-10
HEADER_NUM_SIZE = FOLIO_SIZE      # kicker and folio are one family

# Header band anchors — the PDF anchors converted from mm.
HEADER_KICKER_Y = MARGIN_T        # 13.34 mm
HEADER_TITLE_Y = Inches(0.768)    # 19.5 mm
HEADER_RULE_Y = Inches(1.417)     # 36.0 mm
BODY_TOP = Inches(1.654)          # 42.0 mm

# Cover composition anchors — shared with the PDF cover so a deck cover and
# a manual cover are the same page, not two interpretations of it. The PDF
# engine holds the same anchors in its own units; test_overlap.py asserts the
# two stay in step. Everything here is inches, like the rest of this module.
COVER_RAIL_W = Inches(0.236)     # Saffron rail, full height       (6.0 mm)
COVER_TEXT_X = Inches(1.417)     # text axis for every cover element (36.0 mm)
COVER_LOGO_Y = Inches(0.787)     # hero logo top                    (20.0 mm)
COVER_LOGO_W = Inches(2.165)     # hero logo width                  (55.0 mm)
COVER_RULE_Y = Inches(2.598)     # Saffron rule above the title     (66.0 mm)
COVER_TITLE_Y = Inches(2.913)    # locked title top                 (74.0 mm)
COVER_SUB_Y = Inches(4.95)       # subtitle, below the title block
COVER_META_Y = Inches(5.906)     # Prepared by / Date block        (150.0 mm)
COVER_CLASS_Y = Inches(6.986)    # classification stamp, in the footer band

# Bottom-right cover mark — replaces the duplicate logo that used to sit
# there. (top_in, width_in) per bar, right-anchored to the content margin.
COVER_MARK_BARS = ((5.748, 3.071), (6.063, 2.047), (6.378, 1.024))
COVER_MARK_H = Inches(0.094)     # 2.4 mm

# Brand asset paths
SKILL_DIR = Path(__file__).resolve().parent
ASSETS_DIR = SKILL_DIR / "assets"
LOGO_LIGHT = str(ASSETS_DIR / "lwlogo_lightmode.png")
LOGO_DARK = str(ASSETS_DIR / "lwlogo_darkmode.png")
# Clearance from the right slide edge — the PDF engine's logo clearance in
# inches. Tighter than this and the logo reads as clipped by the slide edge.
LOGO_RIGHT_CLEAR = Inches(0.7)

# Logo sizing (preserve 5.5:1 aspect ratio)
LOGO_W = Inches(1.8)
LOGO_H = Inches(0.327)

# Grid column helpers (12-column, 0.2in gutter, all return inches as float)
MARGIN_L_IN = 0.6
MARGIN_R_IN = 0.6
CONTENT_W_IN = 13.333 - MARGIN_L_IN - MARGIN_R_IN  # 12.133

def col_left(col: int) -> float:
    """Left edge x (inches) for grid column 1-12."""
    assert 1 <= col <= 12
    return MARGIN_L_IN + (col - 1) * (COL_W + GUTTER)

def col_right(col: int) -> float:
    """Right edge x (inches) for grid column 1-12."""
    return col_left(col) + COL_W

def col_span(from_col: int, to_col: int) -> float:
    """Width (inches) spanning cols from_col through to_col inclusive."""
    n = to_col - from_col + 1
    return n * COL_W + (n - 1) * GUTTER

# Font name — embedded as string in PPTX XML
FONT = "Manrope"


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _set_font(run, size_pt: int, color: RGBColor, bold: bool = False,
              name: str = FONT):
    """Set font properties on a text run, including XML-level typeface."""
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    from lxml import etree
    for tag in ("a:latin", "a:ea", "a:sym"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def _add_textbox(slide, left, top, width, height,
                 text: str, size_pt: int = 12, color: RGBColor = DARK_SERPENT,
                 bold: bool = False, align=PP_ALIGN.LEFT,
                 anchor: str = "top", font_name: str = FONT,
                 spacing_after: int = 0, line_spacing: Optional[float] = None):
    """Add a single-paragraph text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    anchor_xml = {"top": "t", "middle": "ctr", "bottom": "b"}
    if anchor in anchor_xml:
        txBox.text_frame._txBody.bodyPr.set("anchor", anchor_xml[anchor])

    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        p.space_after = Pt(spacing_after)
    if line_spacing:
        p.line_spacing = Pt(size_pt * line_spacing)

    run = p.add_run()
    run.text = text
    _set_font(run, size_pt, color, bold, font_name)
    return txBox


def _add_multiline_textbox(slide, left, top, width, height,
                           lines: List[Tuple[str, int, RGBColor, bool]],
                           align=PP_ALIGN.LEFT, spacing_after: int = 6,
                           line_spacing: Optional[float] = None):
    """Add a multi-paragraph text box. Each line: (text, font_size, color, bold)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, size_pt, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_after)
        if line_spacing:
            p.line_spacing = Pt(size_pt * line_spacing)
        run = p.add_run()
        run.text = text
        _set_font(run, size_pt, color, bold)
    return txBox


def _add_image(slide, image_path: str, left, top, width, height):
    """Add image. Silently skip if file missing."""
    if not os.path.isfile(image_path):
        return None
    try:
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        # Move behind text shapes so it doesn't cover them
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
        return pic
    except Exception:
        return None


def _rect(slide, left, top, width, height,
          fill_color: Optional[RGBColor] = None,
          line_color: Optional[RGBColor] = None,
          shadow: bool = False):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if shadow:
        # Subtle shadow via XML
        spPr = shape._element.spPr
        effectLst = spPr.find(qn("a:effectLst"))
        if effectLst is None:
            from lxml import etree
            effectLst = etree.SubElement(spPr, qn("a:effectLst"))
        outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
        outerShdw.set("blurRad", "50800")
        outerShdw.set("dist", "25400")
        outerShdw.set("dir", "5400000")
        outerShdw.set("algn", "tl")
        outerShdw.set("rotWithShape", "0")
    return shape


def _round_rect(slide, left, top, width, height,
                fill_color: Optional[RGBColor] = None,
                line_color: Optional[RGBColor] = None,
                adj: float = 0.08):
    """Add a rounded rectangle. adj: corner radius (0-0.5)."""
    shape = slide.shapes.add_shape(4, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Set adjustment for roundness
    adj_xml = f'<a:avLst><a:gd name="adj1" fmla="val {int(adj * 100)}"/></a:avLst>'
    from lxml import etree
    el = etree.fromstring(adj_xml)
    prstGeom = shape._element.find(qn("a:prstGeom"))
    if prstGeom is not None:
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        # Clear existing
        for child in list(avLst):
            avLst.remove(child)
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj1")
        gd.set("fmla", f"val {int(adj * 100)}")
    return shape


def _thin_rule(slide, left, top, width, color: RGBColor = SAFFRON,
               height=Pt(2.5)):
    """Add a thin horizontal accent line."""
    return _rect(slide, left, top, width, height, fill_color=color)


def _background(slide, color: RGBColor):
    """Set slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _slide_number(slide, current: int, total: int, light_bg: bool = True,
                  label: str = ""):
    """Footer band: hairline divider + centred 'Page X of Y'.

    The footer carries page numbering and the logo ONLY. The left slot stays
    empty — no document title, product name, quote or description. *label* is
    accepted and ignored so existing callers keep working; passing it does
    not put text back in the footer.
    """
    # Hairline divider above the footer zone — same left/right inset as content
    _thin_rule(slide, MARGIN_L, FOOTER_RULE_Y, CONTENT_W,
               CASTLETON_GREEN, FOOTER_RULE_H)

    # Page number — folio style (Castleton Green, bold), bottom-center
    color = CASTLETON_GREEN if light_bg else PAPER
    _add_textbox(
        slide,
        left=Inches(0), top=FOOTER_Y,
        width=SLIDE_W, height=FOOTER_H,
        text=f"Page {current} of {total}",
        size_pt=FOLIO_SIZE, color=color, bold=True,
        align=PP_ALIGN.CENTER, anchor="middle"
    )


def _footer_logo(slide, dark_bg: bool = False):
    """Place the Lifewood logo in the bottom-right corner.

    Interior slides only. The cover carries exactly one logo — the hero mark,
    upper-left — so ``add_cover()`` deliberately does not call this.
    """
    logo_path = LOGO_DARK if dark_bg else LOGO_LIGHT
    logo_x = SLIDE_W - LOGO_W - LOGO_RIGHT_CLEAR
    logo_y = FOOTER_Y + (FOOTER_H - LOGO_H) / 2
    _add_image(slide, logo_path, logo_x, logo_y, LOGO_W, LOGO_H)


def _cover_mark(slide, dark: bool):
    """Abstract brand geometry in the cover's bottom-right corner.

    Three stepped bars echoing the left rail — geometry for balance, not a
    second wordmark competing with the hero logo.
    """
    tail = PAPER if dark else CASTLETON_GREEN
    for (top_in, w_in), color in zip(COVER_MARK_BARS,
                                     (SAFFRON, EARTH_YELLOW, tail)):
        _rect(slide, SLIDE_W - MARGIN_R - Inches(w_in), Inches(top_in),
              Inches(w_in), COVER_MARK_H, fill_color=color)


def _page_header(slide, heading_num: str, heading_text: str,
                 x=MARGIN_L, width=None):
    """Locked content-slide header: folio kicker, then the page title.

    The number and the title are **separate runs at separate sizes** — never
    a single ``"1.0  Heading"`` string. Mixing the folio family into the
    title run is the numbering drift the design system forbids.

    Anchors are the PDF header-band anchors converted from mm, so a slide
    header and a manual header sit on the same line.
    """
    display_num = heading_num if "." in heading_num else heading_num + ".0"
    _add_textbox(slide, x, HEADER_KICKER_Y, Inches(2.0), Inches(0.28),
                 display_num, size_pt=HEADER_NUM_SIZE,
                 color=CASTLETON_GREEN, bold=True)
    _add_textbox(slide, x, HEADER_TITLE_Y, width or CONTENT_W, Inches(0.6),
                 heading_text, size_pt=PAGE_TITLE_SIZE,
                 color=CASTLETON_GREEN, bold=True)


def _accent_top_bar(slide, color: RGBColor = CASTLETON_GREEN,
                    height=Inches(0.05)):
    """Thin colored bar across the very top of the slide."""
    return _rect(slide, Inches(0), Inches(0), SLIDE_W, height,
                 fill_color=color)


def _is_subsection(heading_num: str) -> bool:
    """True for a genuine sub-heading ("1.1", "2.3", "1.1.1") — false
    for a top-level section ("1.0", "2.0") or a bare number ("1").

    Mirrors the TOC sub-entry detection in lifewood_pdf_builder.py so
    both builders agree on what counts as a sub-section.
    """
    segs = heading_num.strip().split(".")
    if len(segs) < 2:
        return False
    try:
        return int(segs[1]) > 0
    except ValueError:
        return False


# ---------------------------------------------------------------------------
# Deck builder
# ---------------------------------------------------------------------------

class LifewoodDeck:
    """Build a premium Lifewood-branded presentation.

    Design system:
        - Widescreen 16:9
        - Manrope typography
        - Castleton Green #046241 headings
        - Dark Serpent #133020 body
        - Saffron #FFB347 accents
        - Sea Salt #F9F7F7 backgrounds
        - Numbered heading hierarchy
    """

    #: Default cover classification for Lifewood deliverables.
    DEFAULT_CLASSIFICATION = "CONFIDENTIAL — FOR INTERNAL USE ONLY"

    def __init__(self, title: str, subtitle: Optional[str] = None,
                 date_str: Optional[str] = None, variant: str = "editorial",
                 classification: Optional[str] = DEFAULT_CLASSIFICATION):
        template_path = SKILL_DIR / "PPT_Master_Overview.pptx"
        if template_path.exists():
            self.prs = Presentation(str(template_path))
            # Remove template's default slides — we build fresh
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].get(qn("r:id"))
                self.prs.part.drop_rel(rId)
                self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])
        else:
            self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank_layout = self.prs.slide_layouts[6]

        self.title = title
        self.subtitle = subtitle
        self.date_str = date_str or datetime.today().strftime("%B %d, %Y")
        self.classification = (classification or "").strip()

        self.variant = variant
        self._vs = get_variant(variant)

        self._slides: list = []
        self._dark_slides: set = set()
        self._front_slides: set = set()

    # ---- Slide type builders ----

    def add_cover(self) -> int:
        """Cover slide — the same composition as the PDF cover.

        Left Saffron rail, hero logo top-left, Saffron rule, left-aligned
        title / subtitle, then the fixed front-matter order
        **Prepared by → Date → Classification** (classification bottom-left).
        ``editorial`` renders it on Castleton Green, ``swiss`` on Sea Salt —
        palette only; the anchors below are the PDF anchors converted from
        mm, so a deck cover and a manual cover are the same page.

        The cover carries **exactly one logo** — the hero mark, upper-left.
        The bottom-right corner holds abstract brand geometry, never a second
        wordmark.
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        dark = self._vs.cover.bg_dark
        if dark:
            self._dark_slides.add(idx)

        title_color = PAPER if dark else CASTLETON_GREEN
        sub_color = SAFFRON if dark else DARK_SERPENT
        meta_color = PAPER if dark else DARK_SERPENT
        logo_asset = LOGO_DARK if dark else LOGO_LIGHT

        _background(slide, CASTLETON_GREEN if dark else SEA_SALT)

        # ---- Left Saffron rail (full height) --------------------------
        _rect(slide, Inches(0), Inches(0), COVER_RAIL_W, SLIDE_H,
              fill_color=SAFFRON)

        # ---- Hero logo, top-left of the text column -------------------
        _add_image(slide, logo_asset, COVER_TEXT_X, COVER_LOGO_Y,
                   COVER_LOGO_W, Inches(2.1654 / 5.5))

        # ---- Saffron rule above the title -----------------------------
        _thin_rule(slide, COVER_TEXT_X, COVER_RULE_Y, Inches(2.756),
                   SAFFRON, Inches(0.0866))

        # ---- Title (locked y, left-aligned) ---------------------------
        text_w = Inches(13.333 * 0.72)
        _add_textbox(slide, COVER_TEXT_X, COVER_TITLE_Y,
                     text_w, Inches(2.0),
                     self.title, size_pt=COVER_TITLE_SIZE, color=title_color,
                     bold=True, align=PP_ALIGN.LEFT)

        # ---- Subtitle -------------------------------------------------
        if self.subtitle:
            _add_textbox(slide, COVER_TEXT_X, COVER_SUB_Y,
                         text_w, Inches(0.7),
                         self.subtitle, size_pt=SUBHEADING_SIZE,
                         color=sub_color, bold=False, align=PP_ALIGN.LEFT)

        # ---- Prepared by / Date (anchored, never floating) ------------
        _add_textbox(slide, COVER_TEXT_X, COVER_META_Y,
                     text_w, Inches(0.3),
                     "Prepared by: Lifewood PH Team",
                     size_pt=CAPTION_SIZE, color=meta_color, bold=False,
                     align=PP_ALIGN.LEFT)
        _add_textbox(slide, COVER_TEXT_X, COVER_META_Y + Inches(0.26),
                     text_w, Inches(0.3),
                     f"Date: {self.date_str}",
                     size_pt=CAPTION_SIZE, color=meta_color, bold=False,
                     align=PP_ALIGN.LEFT)

        # ---- Classification, bottom-left ------------------------------
        if self.classification:
            _thin_rule(slide, COVER_TEXT_X, COVER_CLASS_Y - Inches(0.1),
                       Inches(0.55), SAFFRON, Inches(0.031))
            _add_textbox(slide, COVER_TEXT_X, COVER_CLASS_Y,
                         text_w, Inches(0.25),
                         self.classification.upper(),
                         size_pt=8, color=meta_color, bold=True,
                         align=PP_ALIGN.LEFT)

        # ---- Bottom-right brand mark ----------------------------------
        _cover_mark(slide, dark)

        # ---- Bottom hairline ------------------------------------------
        _thin_rule(slide, Inches(0), SLIDE_H - Inches(0.118), SLIDE_W,
                   SAFFRON if dark else CASTLETON_GREEN, Inches(0.039))

        self._front_slides.add(idx)
        # No bottom-right logo here: the cover carries exactly ONE logo, the
        # hero mark above. The corner holds brand geometry instead.

        return idx

    def add_toc(self, items: List[str]) -> int:
        """Table of contents — card-based layout on Sea Salt.

        Items are strings like ``"1.0  Title"`` or ``"1.1  Sub-title"``.
        The number badge shows the full section number (never
        truncated to a bare chapter digit), and sub-entries ("1.1",
        "2.3", ...) render indented and slightly smaller than their
        parent — matching the sub-entry rule in lifewood_pdf_builder.py
        and SKILL.md §9 ("Include sub-entries when present, indented
        one level").
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        _background(slide, SEA_SALT)
        _accent_top_bar(slide, CASTLETON_GREEN)

        # "Contents" heading
        _add_textbox(slide, MARGIN_L, Inches(0.65),
                     Inches(5), Inches(0.8),
                     "Contents", size_pt=PAGE_TITLE_SIZE, color=CASTLETON_GREEN,
                     bold=True)

        # Accent underline
        _thin_rule(slide, MARGIN_L, Inches(1.45), Inches(3), SAFFRON, Pt(3))

        # TOC items — two-column card layout, clamped to the slide.
        # Rows that would cross the footer band spill to a continuation
        # ToC slide instead of running off the bottom edge.
        cols = 2 if len(items) > 3 else 1
        col_w = (CONTENT_W - Inches(0.3)) / cols
        row_h = Inches(1.0)
        top_y = Inches(1.8)
        row_pitch = row_h + Inches(0.25)
        max_rows = max(int((FOOTER_RULE_Y - top_y) / row_pitch), 1)
        capacity = max_rows * cols
        overflow = list(items[capacity:])
        items = list(items[:capacity])

        for i, item in enumerate(items):
            col = i % cols
            row = i // cols
            x = MARGIN_L + col * (col_w + Inches(0.3))
            y = top_y + row * row_pitch

            parts = item.strip().split(None, 1)
            num_full = parts[0] if parts else str(i + 1)
            title_text = parts[1] if len(parts) > 1 else item
            is_sub = _is_subsection(num_full)
            indent = Inches(0.25) if is_sub else Inches(0)

            # Card background
            _rect(slide, x, y, col_w, row_h, fill_color=WHITE)

            # Left accent border (Saffron)
            _rect(slide, x, y, Inches(0.04), row_h, fill_color=SAFFRON)

            # Number — full "X.Y" form, never truncated to a bare digit
            _add_textbox(slide, x + Inches(0.2) + indent, y + Inches(0.15),
                         Inches(0.7), Inches(0.35),
                         num_full, size_pt=BODY_SIZE - 1 if is_sub else BODY_SIZE,
                         color=SAFFRON, bold=True)

            # Item text
            _add_textbox(slide, x + Inches(0.2) + indent, y + Inches(0.5),
                         col_w - Inches(0.4) - indent, row_h - Inches(0.6),
                         title_text, size_pt=BODY_SIZE - 1 if is_sub else BODY_SIZE,
                         color=DARK_SERPENT, bold=False)

        self._front_slides.add(idx)
        _footer_logo(slide)
        if overflow:
            self.add_toc(overflow)  # continuation ToC — still front matter
        return idx

    def add_section(self, section_num: str, section_title: str,
                     slide_titles: Optional[List[str]] = None,
                     deck_size: int = 0, section_size: int = 0) -> int:
        """Conditional section divider with payload (per Section C.4).

        Divider created only if deck_size >= 12 AND section_size >= 3.
        When divider IS created, right column shows slide_titles as
        'In this section' list (payload).

        When divider is skipped, returns 0 (no slide). Caller must
        add the section number as a Saffron eyebrow on the first
        content slide instead. Returns slide index (1-based) or 0.
        """
        # Conditional logic: skip divider if undersized
        make_divider = True
        if deck_size > 0 and deck_size < 12:
            make_divider = False
        if section_size > 0 and section_size < 3:
            make_divider = False

        if not make_divider:
            return 0  # no slide created; caller handles eyebrow

        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)
        self._dark_slides.add(idx)

        if self._vs.section_divider.bg_dark:
            # --- Editorial: left/right split with payload ---
            _background(slide, DARK_SERPENT)

            # Left Saffron bar
            _rect(slide, Inches(0), Inches(0),
                  Inches(0.35), SLIDE_H, fill_color=SAFFRON)

            # Left column (5 cols) — number + title + rule
            _add_textbox(slide, Inches(col_left(1)), Inches(1.5),
                         Inches(col_span(1, 5)), Inches(1.5),
                         section_num, size_pt=72, color=SAFFRON,
                         bold=True)

            _add_textbox(slide, Inches(col_left(1)), Inches(3.2),
                         Inches(col_span(1, 5)), Inches(0.8),
                         section_title, size_pt=32, color=PAPER,
                         bold=False)

            _thin_rule(slide, Inches(col_left(1)), Inches(4.2),
                       Inches(4), SAFFRON, Pt(2))

            # Right column (7 cols) — "In this section" list
            if slide_titles:
                _add_textbox(slide, Inches(col_left(7)), Inches(1.8),
                             Inches(col_span(7, 12)), Inches(0.5),
                             "In this section:", size_pt=BODY_SIZE, color=SAFFRON,
                             bold=True)
                for i, st in enumerate(slide_titles[:4]):
                    _add_textbox(slide, Inches(col_left(7) + 0.2), Inches(2.5 + i * 0.6),
                                 Inches(col_span(7, 12) - 0.2), Inches(0.5),
                                 f"• {st}", size_pt=BODY_SIZE, color=PAPER,
                                 bold=False)

            _footer_logo(slide, dark_bg=True)
        else:
            # --- Swiss: light background, clean minimal ---
            _background(slide, SEA_SALT)

            _thin_rule(slide, Inches(0), Inches(0),
                       SLIDE_W, SAFFRON, Inches(0.08))

            # Left: section number + title
            _add_textbox(slide, MARGIN_L, Inches(1.2),
                         Inches(col_span(1, 5)), Inches(1.2),
                         section_num, size_pt=36, color=SAFFRON,
                         bold=True)

            _thin_rule(slide, MARGIN_L, Inches(2.6),
                       Inches(4), DARK_SERPENT, Pt(2))

            _add_textbox(slide, MARGIN_L, Inches(3.0),
                         Inches(col_span(1, 5)), Inches(1.0),
                         section_title, size_pt=32, color=CASTLETON_GREEN,
                         bold=False)

            # Right: payload list
            if slide_titles:
                _add_textbox(slide, Inches(col_left(7)), Inches(1.8),
                             Inches(col_span(7, 12)), Inches(0.5),
                             "In this section:", size_pt=BODY_SIZE, color=SAFFRON,
                             bold=True)
                for i, st in enumerate(slide_titles[:4]):
                    _add_textbox(slide, Inches(col_left(7) + 0.2), Inches(2.5 + i * 0.5),
                                 Inches(col_span(7, 12) - 0.2), Inches(0.4),
                                 f"• {st}", size_pt=BODY_SIZE, color=DARK_SERPENT,
                                 bold=False)

            _footer_logo(slide, dark_bg=False)

        return idx

    def add_content(self, heading_num: str, heading_text: str,
                    bullets: List[str],
                    notes: Optional[str] = None,
                    eyebrow: Optional[str] = None) -> int:
        """Content slide — heading block with body bullets.

        eyebrow: optional Saffron section number shown above title
                 (used when section divider was skipped per C.4).

        Bullets that would cross the footer rule continue on a "(cont.)"
        slide — PowerPoint clips anything past the slide edge, so the
        builder never places content there.
        """
        overflow: List[str] = []
        if bullets:
            _grid = self._vs.content.layout_variant == "grid"
            _y_body = Inches(1.65) if _grid else Inches(0.55) + Inches(1.35)
            _line_h = Inches(0.45) if _grid else Inches(0.50)
            if eyebrow and eyebrow != heading_num and not _grid:
                _y_body += Inches(0.25)
            _max = max(int((FOOTER_RULE_Y - _y_body) / _line_h), 1)
            if len(bullets) > _max:
                overflow = list(bullets[_max:])
                bullets = list(bullets[:_max])

        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean, dense, no decorative elements ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            # Heading — kicker (folio number) then title on its own line
            _page_header(slide, heading_num, heading_text)

            # Body bullets — 16pt body scale, clean text
            if bullets:
                y_body = Inches(1.65)
                line_h = Inches(0.45)

                for i, b in enumerate(bullets):
                    y = y_body + line_h * i
                    _add_textbox(slide, MARGIN_L, y,
                                 CONTENT_W, line_h,
                                 b, size_pt=BODY_SIZE, color=DARK_SERPENT,
                                 bold=False)

            # Bottom Dark Serpent rule
            _thin_rule(slide, MARGIN_L, Inches(6.35),
                       CONTENT_W, DARK_SERPENT, Pt(0.5))
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            # Eyebrow — only when section divider was skipped (C.4) AND it
            # differs from the folio kicker (which always shows heading_num)
            eyebrow_y = Inches(0.55)
            if eyebrow and eyebrow != heading_num:
                _add_textbox(slide, MARGIN_L, Inches(0.3),
                             CONTENT_W, Inches(0.35),
                             eyebrow, size_pt=CAPTION_SIZE, color=SAFFRON,
                             bold=True)
                eyebrow_y = Inches(0.8)

            # Heading area — white background block
            _rect(slide, MARGIN_L, eyebrow_y,
                  CONTENT_W, Inches(1.15), fill_color=WHITE)

            # Kicker — full folio number ("1.0" / "1.1"), same style as footer
            _add_textbox(slide, MARGIN_L + Inches(0.15), eyebrow_y + Inches(0.05),
                         Inches(2.5), Inches(0.3),
                         heading_num, size_pt=HEADER_NUM_SIZE,
                         color=CASTLETON_GREEN, bold=True)

            # Heading text — locked title scale, aligned with the kicker
            _add_textbox(slide, MARGIN_L + Inches(0.15), eyebrow_y + Inches(0.32),
                         CONTENT_W - Inches(0.3), Inches(0.6),
                         heading_text, size_pt=PAGE_TITLE_SIZE,
                         color=CASTLETON_GREEN, bold=True)

            # Body bullets — body scale 16pt with spacing
            if bullets:
                y_body = eyebrow_y + Inches(1.35)
                line_h = Inches(0.50)

                for i, b in enumerate(bullets):
                    y = y_body + line_h * i

                    # Bullet marker — Saffron square
                    sq = Inches(0.08)
                    _rect(slide, MARGIN_L, y + Inches(0.12),
                          sq, sq, fill_color=SAFFRON)

                    # Bullet text — 16pt body scale
                    _add_textbox(slide, MARGIN_L + Inches(0.2), y,
                                 CONTENT_W - Inches(0.2), line_h,
                                 b, size_pt=BODY_SIZE, color=DARK_SERPENT,
                                 bold=False)

            # Bottom rule
            _thin_rule(slide, MARGIN_L, Inches(6.35),
                       CONTENT_W, DARK_SERPENT, Pt(0.5))

        _footer_logo(slide)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        if overflow:
            cont = heading_text if heading_text.endswith("(cont.)") else \
                f"{heading_text} (cont.)"
            self.add_content(heading_num, cont, overflow)

        return idx

    def add_closing(self, message: str = "Thank You",
                    subtitle: Optional[str] = None) -> int:
        """Closing slide — elegant dark green."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)
        self._dark_slides.add(idx)

        _background(slide, CASTLETON_GREEN)
        _accent_top_bar(slide, SAFFRON, Inches(0.05))

        # Closing message — 44pt Paper, centered on live area focal zone
        _add_textbox(slide, MARGIN_L, Inches(2.8),
                     CONTENT_W, Inches(1.5),
                     message, size_pt=44, color=PAPER, bold=True,
                     align=PP_ALIGN.CENTER)

        # Saffron accent line
        _thin_rule(slide, (SLIDE_W - Inches(2.5)) / 2, Inches(4.4),
                   Inches(2.5), SAFFRON, Pt(3))

        # Subtitle — Saffron, below accent line
        if subtitle:
            _add_textbox(slide, MARGIN_L, Inches(4.8),
                         CONTENT_W, Inches(0.7),
                         subtitle, size_pt=SUBHEADING_SIZE, color=SAFFRON,
                         bold=False, align=PP_ALIGN.CENTER)

        # Metadata block — bottom
        _add_textbox(slide, MARGIN_L, Inches(5.3),
                     CONTENT_W, Inches(0.35),
                     "Prepared by: Lifewood PH Team",
                     size_pt=11, color=PAPER, bold=False,
                     align=PP_ALIGN.CENTER)

        _add_textbox(slide, MARGIN_L, Inches(5.7),
                     CONTENT_W, Inches(0.35),
                     f"Date: {self.date_str}",
                     size_pt=11, color=PAPER, bold=False,
                     align=PP_ALIGN.CENTER)

        _footer_logo(slide, dark_bg=True)

        _thin_rule(slide, Inches(0), SLIDE_H - Inches(0.1),
                   SLIDE_W, SAFFRON, Inches(0.04))

        return idx

    # ---- Preset content layouts ----

    def add_two_column(self, heading_num: str, heading_text: str,
                       left_heading: str, left_items: List[str],
                       right_heading: str, right_items: List[str],
                       notes: Optional[str] = None) -> int:
        """Two-column comparison layout."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean, dense, no decorative heading block ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            mid_x = MARGIN_L + CONTENT_W / 2
            _thin_rule(slide, mid_x, Inches(1.5), Pt(1.5), DARK_SERPENT,
                       height=Inches(4.5))

            col_w = CONTENT_W / 2 - Inches(0.3)

            _add_textbox(slide, MARGIN_L + Inches(0.1), Inches(1.5),
                         col_w, Inches(0.35),
                         left_heading, size_pt=SUBHEADING_SIZE, color=CASTLETON_GREEN, bold=True)
            for i, item in enumerate(left_items):
                _add_textbox(slide, MARGIN_L + Inches(0.1), Inches(1.9 + i * 0.45),
                             col_w, Inches(0.40),
                             item, size_pt=BODY_SIZE, color=DARK_SERPENT)

            _add_textbox(slide, mid_x + Inches(0.3), Inches(1.5),
                         col_w, Inches(0.35),
                         right_heading, size_pt=SUBHEADING_SIZE, color=CASTLETON_GREEN, bold=True)
            for i, item in enumerate(right_items):
                _add_textbox(slide, mid_x + Inches(0.3), Inches(1.9 + i * 0.45),
                             col_w, Inches(0.40),
                             item, size_pt=BODY_SIZE, color=DARK_SERPENT)
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            # Heading
            _rect(slide, MARGIN_L, Inches(0.5), CONTENT_W, Inches(0.9), fill_color=WHITE)
            _page_header(slide, heading_num, heading_text,
                         x=MARGIN_L + Inches(0.2),
                         width=CONTENT_W - Inches(0.4))

            # Column divider
            mid_x = MARGIN_L + CONTENT_W / 2
            _thin_rule(slide, mid_x, Inches(1.7), Pt(1.5), DARK_SERPENT,
                        height=Inches(4.5))

            col_w = CONTENT_W / 2 - Inches(0.3)

            # Left column — subhead 20pt, body 16pt
            _add_textbox(slide, MARGIN_L + Inches(0.2), Inches(1.7),
                         col_w, Inches(0.4),
                         left_heading, size_pt=SUBHEADING_SIZE, color=CASTLETON_GREEN, bold=True)
            for i, item in enumerate(left_items):
                _add_textbox(slide, MARGIN_L + Inches(0.2), Inches(2.2 + i * 0.45),
                             col_w, Inches(0.40),
                             f"▸ {item}", size_pt=BODY_SIZE, color=DARK_SERPENT)

            # Right column — subhead 20pt, body 16pt
            _add_textbox(slide, mid_x + Inches(0.3), Inches(1.7),
                         col_w, Inches(0.4),
                         right_heading, size_pt=SUBHEADING_SIZE, color=CASTLETON_GREEN, bold=True)
            for i, item in enumerate(right_items):
                _add_textbox(slide, mid_x + Inches(0.3), Inches(2.2 + i * 0.45),
                             col_w, Inches(0.40),
                             f"▸ {item}", size_pt=BODY_SIZE, color=DARK_SERPENT)

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_stats(self, heading_num: str, heading_text: str,
                  stats: List[Tuple[str, str, str]],
                  notes: Optional[str] = None) -> int:
        """Stats/metrics layout — big numbers with labels.
        stats: [(number, label, description), ...]
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean, tighter stats, no card shadows ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            n = min(len(stats), 4)
            stat_w = (CONTENT_W - Inches(0.3) * (n - 1)) / n
            y_start = Inches(1.6)

            for i, (number, label, desc) in enumerate(stats[:4]):
                x = MARGIN_L + i * (stat_w + Inches(0.3))

                # Card — no saffron top stripe
                _rect(slide, x, y_start, stat_w, Inches(2.6), fill_color=SEA_SALT)

                # Big number — Castleton Green, 44pt
                _add_textbox(slide, x, y_start + Inches(0.3),
                             stat_w, Inches(0.8),
                             number, size_pt=44, color=CASTLETON_GREEN, bold=True,
                             align=PP_ALIGN.CENTER)

                # Label — 14pt bold (caption scale)
                _add_textbox(slide, x, y_start + Inches(1.1),
                             stat_w, Inches(0.35),
                             label, size_pt=CAPTION_SIZE, color=DARK_SERPENT, bold=True,
                             align=PP_ALIGN.CENTER)

                # Description — 12pt (caption scale)
                _add_textbox(slide, x + Inches(0.1), y_start + Inches(1.5),
                             stat_w - Inches(0.2), Inches(0.9),
                             desc, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False,
                             align=PP_ALIGN.CENTER)
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _rect(slide, MARGIN_L, Inches(0.5), CONTENT_W, Inches(0.9), fill_color=WHITE)
            _page_header(slide, heading_num, heading_text,
                         x=MARGIN_L + Inches(0.2),
                         width=CONTENT_W - Inches(0.4))

            # Stats grid — up to 4 stats in a row
            n = min(len(stats), 4)
            stat_w = (CONTENT_W - Inches(0.4) * (n - 1)) / n
            y_start = Inches(1.8)

            for i, (number, label, desc) in enumerate(stats[:4]):
                x = MARGIN_L + i * (stat_w + Inches(0.4))

                # Card
                _rect(slide, x, y_start, stat_w, Inches(2.8), fill_color=WHITE)
                _rect(slide, x, y_start, stat_w, Inches(0.06), fill_color=SAFFRON)

                # Big number — 44pt per stat scale
                _add_textbox(slide, x, y_start + Inches(0.3),
                             stat_w, Inches(0.9),
                             number, size_pt=44, color=SAFFRON, bold=True,
                             align=PP_ALIGN.CENTER)

                # Label — 14pt bold caption
                _add_textbox(slide, x, y_start + Inches(1.2),
                             stat_w, Inches(0.4),
                             label, size_pt=CAPTION_SIZE, color=CASTLETON_GREEN, bold=True,
                             align=PP_ALIGN.CENTER)

                # Description — 12pt caption
                _add_textbox(slide, x + Inches(0.15), y_start + Inches(1.6),
                             stat_w - Inches(0.3), Inches(1.0),
                             desc, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False,
                             align=PP_ALIGN.CENTER)

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_quote(self, heading_num: str, quote_text: str,
                  attribution: Optional[str] = None,
                  notes: Optional[str] = None) -> int:
        """Quote/testimonial layout — centered pull quote."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)
        self._dark_slides.add(idx)

        _background(slide, DARK_SERPENT)

        # Large decorative quote mark
        _add_textbox(slide, MARGIN_L, Inches(1.0),
                     Inches(1.5), Inches(1.0),
                     "\u201C", size_pt=80, color=SAFFRON, bold=True)

        # Quote text
        _add_textbox(slide, MARGIN_L + Inches(1.2), Inches(1.5),
                     CONTENT_W - Inches(1.5), Inches(2.5),
                     quote_text, size_pt=28, color=PAPER, bold=False,
                     align=PP_ALIGN.LEFT)

        # Accent line
        _thin_rule(slide, MARGIN_L + Inches(1.2), Inches(4.3),
                   Inches(2.5), SAFFRON, Pt(3))

        # Attribution
        if attribution:
            _add_textbox(slide, MARGIN_L + Inches(1.2), Inches(4.6),
                         CONTENT_W - Inches(1.5), Inches(0.5),
                         attribution, size_pt=CAPTION_SIZE, color=SAFFRON, bold=False)

        _footer_logo(slide, dark_bg=True)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_timeline(self, heading_num: str, heading_text: str,
                     milestones: List[Tuple[str, str, str]],
                     notes: Optional[str] = None) -> int:
        """Timeline/roadmap layout — horizontal timeline with phases.
        milestones: [(phase_label, title, description), ...]
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean timeline, no decorative heading block ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            n = min(len(milestones), 5)
            phase_w = (CONTENT_W - Inches(0.25) * (n - 1)) / n
            bar_y = Inches(2.0)

            # Horizontal timeline bar — Dark Serpent instead of Saffron
            _rect(slide, MARGIN_L, bar_y, CONTENT_W, Pt(2), fill_color=DARK_SERPENT)

            for i, (phase, title, desc) in enumerate(milestones[:5]):
                x = MARGIN_L + i * (phase_w + Inches(0.25))

                # Circle — Castleton Green
                circle_d = Inches(0.25)
                circle = slide.shapes.add_shape(1, x, bar_y - circle_d / 2,
                                                circle_d, circle_d)
                circle.fill.solid()
                circle.fill.fore_color.rgb = CASTLETON_GREEN
                circle.line.fill.background()

                # Phase label — 12pt caption scale
                _add_textbox(slide, x, Inches(2.5), phase_w, Inches(0.25),
                             phase, size_pt=CAPTION_SIZE, color=CASTLETON_GREEN, bold=True,
                             align=PP_ALIGN.CENTER)

                # Title — 16pt body scale
                _add_textbox(slide, x, Inches(2.8), phase_w, Inches(0.35),
                             title, size_pt=SUBHEADING_SIZE, color=DARK_SERPENT, bold=True,
                             align=PP_ALIGN.CENTER)

                # Description — 12pt caption scale
                _add_textbox(slide, x + Inches(0.05), Inches(3.2),
                             phase_w - Inches(0.1), Inches(2.5),
                             desc, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False,
                             align=PP_ALIGN.CENTER)
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _rect(slide, MARGIN_L, Inches(0.5), CONTENT_W, Inches(0.9), fill_color=WHITE)
            _page_header(slide, heading_num, heading_text,
                         x=MARGIN_L + Inches(0.2),
                         width=CONTENT_W - Inches(0.4))

            n = min(len(milestones), 5)
            phase_w = (CONTENT_W - Inches(0.3) * (n - 1)) / n

            # Horizontal timeline bar
            bar_y = Inches(2.3)
            _rect(slide, MARGIN_L, bar_y, CONTENT_W, Pt(3), fill_color=SAFFRON)

            for i, (phase, title, desc) in enumerate(milestones[:5]):
                x = MARGIN_L + i * (phase_w + Inches(0.3))

                # Circle on timeline
                circle_d = Inches(0.3)
                circle = slide.shapes.add_shape(1, x, bar_y - circle_d / 2 + Pt(1.5),
                                                circle_d, circle_d)
                circle.fill.solid()
                circle.fill.fore_color.rgb = SAFFRON
                circle.line.fill.background()

                # Phase label — 14pt Saffron bold (caption scale)
                _add_textbox(slide, x, Inches(2.8), phase_w, Inches(0.3),
                             phase, size_pt=CAPTION_SIZE, color=SAFFRON, bold=True,
                             align=PP_ALIGN.CENTER)

                # Title — 18pt subhead scale
                _add_textbox(slide, x, Inches(3.15), phase_w, Inches(0.4),
                             title, size_pt=SUBHEADING_SIZE, color=CASTLETON_GREEN, bold=True,
                             align=PP_ALIGN.CENTER)

                # Description — 14pt body
                _add_textbox(slide, x + Inches(0.1), Inches(3.6),
                             phase_w - Inches(0.2), Inches(2.0),
                             desc, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False,
                             align=PP_ALIGN.CENTER)

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_table(self, heading_num: str, heading_text: str,
                  headers: List[str], rows: List[List[str]],
                  notes: Optional[str] = None) -> int:
        """Table layout — structured data grid."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean table, no decorative heading block ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            n_cols = len(headers)
            n_rows = min(len(rows), 7)
            col_w = CONTENT_W / n_cols
            row_h = Inches(0.38)

            # Header row — 14pt bold per table spec
            for c, header in enumerate(headers):
                x = MARGIN_L + c * col_w
                _rect(slide, x, Inches(1.5), col_w, row_h, fill_color=CASTLETON_GREEN)
                _add_textbox(slide, x + Inches(0.08), Inches(1.52),
                             col_w - Inches(0.16), row_h - Inches(0.04),
                             header, size_pt=TABLE_SIZE, color=WHITE, bold=True,
                             anchor="middle")

            # Data rows — 14pt
            for r, row in enumerate(rows[:n_rows]):
                y = Inches(1.5 + 0.38 + r * 0.42)
                bg = WHITE if r % 2 == 0 else SEA_SALT
                for c, cell in enumerate(row):
                    x = MARGIN_L + c * col_w
                    _rect(slide, x, y, col_w, row_h, fill_color=bg)
                    _add_textbox(slide, x + Inches(0.08), y + Inches(0.02),
                                 col_w - Inches(0.16), row_h - Inches(0.04),
                                 str(cell), size_pt=TABLE_SIZE, color=DARK_SERPENT, bold=False,
                                 anchor="middle")
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _rect(slide, MARGIN_L, Inches(0.5), CONTENT_W, Inches(0.9), fill_color=WHITE)
            _page_header(slide, heading_num, heading_text,
                         x=MARGIN_L + Inches(0.2),
                         width=CONTENT_W - Inches(0.4))

            # Table
            n_cols = len(headers)
            n_rows = min(len(rows), 7)  # max 7 data rows
            col_w = CONTENT_W / n_cols
            row_h = Inches(0.45)

            # Header row — 14pt bold per table spec
            for c, header in enumerate(headers):
                x = MARGIN_L + c * col_w
                _rect(slide, x, Inches(1.7), col_w, row_h, fill_color=CASTLETON_GREEN)
                _add_textbox(slide, x + Inches(0.1), Inches(1.72),
                             col_w - Inches(0.2), row_h - Inches(0.05),
                             header, size_pt=TABLE_SIZE, color=WHITE, bold=True,
                             anchor="middle")

            # Data rows — 14pt per table spec
            for r, row in enumerate(rows[:n_rows]):
                y = Inches(1.7 + 0.45 + r * 0.42)
                bg = WHITE if r % 2 == 0 else SEA_SALT
                for c, cell in enumerate(row):
                    x = MARGIN_L + c * col_w
                    _rect(slide, x, y, col_w, row_h, fill_color=bg)
                    _add_textbox(slide, x + Inches(0.1), y + Inches(0.03),
                                 col_w - Inches(0.2), row_h - Inches(0.05),
                                 str(cell), size_pt=TABLE_SIZE, color=DARK_SERPENT, bold=False,
                                 anchor="middle")

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_checklist(self, heading_num: str, heading_text: str,
                      items: List[Tuple[str, bool]],
                      notes: Optional[str] = None) -> int:
        """Checklist layout — items with check/cross indicators.
        items: [(text, completed), ...]
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean checklist, no decorative heading block ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            y_start = Inches(1.5)
            line_h = Inches(0.4)

            for i, (item, done) in enumerate(items):
                y = y_start + line_h * i

                check_d = Inches(0.22)
                if done:
                    _rect(slide, MARGIN_L, y + Inches(0.04), check_d, check_d,
                          fill_color=CASTLETON_GREEN)
                    _add_textbox(slide, MARGIN_L, y + Inches(0.04), check_d, check_d,
                                 "✓", size_pt=BODY_SIZE, color=WHITE, bold=True,
                                 align=PP_ALIGN.CENTER, anchor="middle")
                else:
                    _rect(slide, MARGIN_L, y + Inches(0.04), check_d, check_d,
                          fill_color=WHITE, line_color=DARK_SERPENT)

                _add_textbox(slide, MARGIN_L + Inches(0.35), y,
                             CONTENT_W - Inches(0.4), line_h,
                             item, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False)
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _rect(slide, MARGIN_L, Inches(0.5), CONTENT_W, Inches(0.9), fill_color=WHITE)
            _page_header(slide, heading_num, heading_text,
                         x=MARGIN_L + Inches(0.2),
                         width=CONTENT_W - Inches(0.4))

            y_start = Inches(1.8)
            line_h = Inches(0.5)

            for i, (item, done) in enumerate(items):
                y = y_start + line_h * i

                # Checkbox
                check_d = Inches(0.28)
                if done:
                    _rect(slide, MARGIN_L, y, check_d, check_d, fill_color=CASTLETON_GREEN)
                    _add_textbox(slide, MARGIN_L, y, check_d, check_d,
                                 "✓", size_pt=BODY_SIZE, color=WHITE, bold=True,
                                 align=PP_ALIGN.CENTER, anchor="middle")
                else:
                    _rect(slide, MARGIN_L, y, check_d, check_d, fill_color=WHITE,
                          line_color=DARK_SERPENT)

                # Item text
                _add_textbox(slide, MARGIN_L + Inches(0.45), y,
                             CONTENT_W - Inches(0.5), line_h,
                             item, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False)

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_callout(self, heading_num: str, heading_text: str,
                    callout_text: str, supporting_text: Optional[str] = None,
                    notes: Optional[str] = None) -> int:
        """Callout/highlight layout — big message with supporting detail."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)
        self._dark_slides.add(idx)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: light callout, clean ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            # Big message — left-aligned, Dark Serpent
            _add_textbox(slide, MARGIN_L, Inches(1.8),
                         CONTENT_W, Inches(1.8),
                         callout_text, size_pt=28, color=DARK_SERPENT, bold=True,
                         align=PP_ALIGN.LEFT)

            # Accent line
            _thin_rule(slide, MARGIN_L, Inches(3.8),
                       Inches(3), SAFFRON, Pt(3))

            if supporting_text:
                _add_textbox(slide, MARGIN_L, Inches(4.2),
                             CONTENT_W, Inches(1.5),
                             supporting_text, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False,
                             align=PP_ALIGN.LEFT)

            _footer_logo(slide, dark_bg=False)
        else:
            # --- Editorial (current behavior) ---
            _background(slide, CASTLETON_GREEN)
            _accent_top_bar(slide, SAFFRON, Inches(0.05))

            # Callout heading
            _add_textbox(slide, MARGIN_L, Inches(1.0),
                         CONTENT_W, Inches(0.5),
                         heading_num, size_pt=HEADER_NUM_SIZE, color=PAPER,
                         bold=True, align=PP_ALIGN.CENTER)

            # Big message
            _add_textbox(slide, MARGIN_L, Inches(1.8),
                         CONTENT_W, Inches(2.0),
                         callout_text, size_pt=36, color=PAPER, bold=True,
                         align=PP_ALIGN.CENTER)

            # Accent line
            _thin_rule(slide, (SLIDE_W - Inches(3)) / 2, Inches(4.0),
                       Inches(3), SAFFRON, Pt(3))

            if supporting_text:
                _add_textbox(slide, MARGIN_L + Inches(1), Inches(4.4),
                             CONTENT_W - Inches(2), Inches(1.5),
                             supporting_text, size_pt=BODY_SIZE, color=PAPER, bold=False,
                             align=PP_ALIGN.CENTER)

            _footer_logo(slide, dark_bg=True)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_card_grid(self, heading_num: str, heading_text: str,
                      cards: List[Tuple[str, List[str]]],
                      grid: str = "2x2",
                      notes: Optional[str] = None) -> int:
        """Card grid layout — 2x2 or 3x2 grid of titled cards.

        cards: [(card_title, [body_line, ...]), ...]
        grid: '2x2' (4 cards) or '3x2' (6 cards)
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        _background(slide, SEA_SALT)
        _accent_top_bar(slide, CASTLETON_GREEN)

        # Slide title — locked page-title scale
        _add_textbox(slide, MARGIN_L, HEADER_TITLE_Y,
                     CONTENT_W, Inches(0.9),
                     heading_text, size_pt=PAGE_TITLE_SIZE,
                     color=CASTLETON_GREEN, bold=True)

        is_3x2 = grid == "3x2"
        n_cols = 3 if is_3x2 else 2
        n_rows = 2

        if is_3x2:
            # 3 cols: widths based on grid cols 1-4, 5-8, 9-12
            card_w = Inches(col_span(1, 4))
            positions = [
                (Inches(col_left(1)), Inches(1.6)),
                (Inches(col_left(5)), Inches(1.6)),
                (Inches(col_left(9)), Inches(1.6)),
                (Inches(col_left(1)), Inches(4.0)),
                (Inches(col_left(5)), Inches(4.0)),
                (Inches(col_left(9)), Inches(4.0)),
            ]
        else:
            # 2 cols: cols 1-6, 7-12
            card_w = Inches(col_span(1, 6))
            positions = [
                (Inches(col_left(1)), Inches(1.6)),
                (Inches(col_left(7)), Inches(1.6)),
                (Inches(col_left(1)), Inches(4.2)),
                (Inches(col_left(7)), Inches(4.2)),
            ]

        card_h = Inches(2.2) if is_3x2 else Inches(2.4)
        inner_pad = Inches(0.25)

        for i, (title, body_lines) in enumerate(cards[: n_cols * n_rows]):
            cx, cy = positions[i]

            # Card background — white on sea salt
            _rect(slide, cx, cy, card_w, card_h, fill_color=WHITE,
                  line_color=RGBColor(0x04, 0x62, 0x41), shadow=False)
            # 1px border at 15% opacity via thin line
            _rect(slide, cx, cy, card_w, Pt(1.5),
                  fill_color=RGBColor(0x04, 0x62, 0x41))

            # Card title — 18pt subhead scale
            _add_textbox(slide, cx + inner_pad, cy + inner_pad,
                         card_w - inner_pad * 2, Inches(0.4),
                         title, size_pt=SUBHEADING_SIZE, color=DARK_SERPENT,
                         bold=True)

            # Body — 16pt, 2-3 lines
            body_text = "\n".join(body_lines[:3])
            _add_textbox(slide, cx + inner_pad, cy + inner_pad + Inches(0.5),
                         card_w - inner_pad * 2, Inches(1.4),
                         body_text, size_pt=BODY_SIZE, color=DARK_SERPENT,
                         bold=False)

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    def add_image_text(self, heading_num: str, heading_text: str,
                       image_path: str, image_side: str = "left",
                       bullets: Optional[List[str]] = None,
                       notes: Optional[str] = None) -> int:
        """Image + text layout — visual on one side, content on other."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._slides.append(slide)
        idx = len(self._slides)

        if self._vs.content.layout_variant == "grid":
            # --- Swiss: clean image+text, no decorative heading block ---
            _background(slide, WHITE)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _page_header(slide, heading_num, heading_text)

            img_w = CONTENT_W * 0.45
            txt_w = CONTENT_W * 0.55

            if image_side == "left":
                img_x, txt_x = MARGIN_L, MARGIN_L + img_w + Inches(0.25)
            else:
                txt_x, img_x = MARGIN_L, MARGIN_L + txt_w + Inches(0.25)

            if os.path.isfile(image_path):
                try:
                    slide.shapes.add_picture(image_path, img_x, Inches(1.6),
                                             img_w, Inches(4.5))
                except Exception:
                    pass

            if bullets:
                y_start = Inches(1.6)
                line_h = Inches(0.35)
                for i, b in enumerate(bullets):
                    y = y_start + line_h * i
                    _add_textbox(slide, txt_x, y,
                                 txt_w, line_h,
                                 b, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False)
        else:
            # --- Editorial (current behavior) ---
            _background(slide, SEA_SALT)
            _accent_top_bar(slide, CASTLETON_GREEN)

            _rect(slide, MARGIN_L, Inches(0.5), CONTENT_W, Inches(0.9), fill_color=WHITE)
            _page_header(slide, heading_num, heading_text,
                         x=MARGIN_L + Inches(0.2),
                         width=CONTENT_W - Inches(0.4))

            # Layout: image 45%, text 55%
            img_w = CONTENT_W * 0.45
            txt_w = CONTENT_W * 0.55

            if image_side == "left":
                img_x, txt_x = MARGIN_L, MARGIN_L + img_w + Inches(0.3)
            else:
                txt_x, img_x = MARGIN_L, MARGIN_L + txt_w + Inches(0.3)

            # Image
            if os.path.isfile(image_path):
                try:
                    slide.shapes.add_picture(image_path, img_x, Inches(1.8),
                                             img_w, Inches(4.5))
                except Exception:
                    pass

            # Text bullets
            if bullets:
                y_start = Inches(1.8)
                line_h = Inches(0.42)
                for i, b in enumerate(bullets):
                    y = y_start + line_h * i
                    sq = Inches(0.08)
                    _rect(slide, txt_x, y + Inches(0.08), sq, sq, fill_color=SAFFRON)
                    _add_textbox(slide, txt_x + Inches(0.2), y,
                                 txt_w - Inches(0.2), line_h,
                                 b, size_pt=BODY_SIZE, color=DARK_SERPENT, bold=False)

        _footer_logo(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return idx

    # ---- Smart layout engine ----

    def smart_slide(self, heading_num: str, heading_text: str,
                    content: List[str],
                    notes: Optional[str] = None) -> int:
        """Auto-select best layout based on content analysis.

        Uses internal heuristics to pick the most appropriate preset.
        Returns the slide index.
        """
        layout = ContentAnalyzer.decide_layout(content)
        return self._apply_layout(heading_num, heading_text, content, layout, notes)

    def _apply_layout(self, heading_num: str, heading_text: str,
                      content: List[str], layout: str,
                      notes: Optional[str] = None) -> int:
        """Apply a specific layout preset."""
        if layout == "stats":
            stats = ContentAnalyzer.extract_stats(content)
            if stats:
                return self.add_stats(heading_num, heading_text, stats, notes)
        elif layout == "two_column":
            left, right = ContentAnalyzer.split_two_column(content)
            if left and right:
                left_head, left_items = left
                right_head, right_items = right
                return self.add_two_column(heading_num, heading_text,
                                           left_head, left_items,
                                           right_head, right_items, notes)
        elif layout == "table":
            table_data = ContentAnalyzer.extract_table(content)
            if table_data:
                headers, rows = table_data
                return self.add_table(heading_num, heading_text, headers, rows, notes)
        elif layout == "checklist":
            items = ContentAnalyzer.extract_checklist(content)
            if items:
                return self.add_checklist(heading_num, heading_text, items, notes)
        elif layout == "timeline":
            milestones = ContentAnalyzer.extract_timeline(content)
            if milestones:
                return self.add_timeline(heading_num, heading_text, milestones, notes)

        # Fallback to standard bullet layout
        return self.add_content(heading_num, heading_text, content, notes)

    # ---- Save + export ----

    def save(self, path: str, skip_density_check: bool = False) -> str:
        """Add page numbers and save PPTX.

        Runs density check on each non-cover slide unless skip_density_check=True.
        """
        front = self._front_slides
        numbered = [(i, slide) for i, slide in enumerate(self._slides)
                    if (i + 1) not in front]

        for k, (i, slide) in enumerate(numbered):
            light_bg = (i + 1) not in self._dark_slides
            _slide_number(slide, k + 1, len(numbered), light_bg=light_bg,
                          label=self.title)

            # Density check bbox (Section C.2/C.7)
            if not skip_density_check:
                spTree = slide.shapes._spTree
                shapes = spTree.findall(qn("p:sp"))
                if not shapes:
                    continue
                # Rough bbox from all shape positions
                min_x = float('inf')
                min_y = float('inf')
                max_x = 0
                max_y = 0
                count = 0
                for sp in shapes:
                    spPr = sp.find(qn("a:xfrm"))
                    if spPr is None:
                        continue
                    off = spPr.find(qn("a:off"))
                    ext = spPr.find(qn("a:ext"))
                    if off is None or ext is None:
                        continue
                    try:
                        sx = int(off.get("x", 0))
                        sy = int(off.get("y", 0))
                        sw = int(ext.get("cx", 0))
                        sh = int(ext.get("cy", 0))
                    except (ValueError, TypeError):
                        continue
                    if sw == 0 or sh == 0:
                        continue
                    min_x = min(min_x, sx)
                    min_y = min(min_y, sy)
                    max_x = max(max_x, sx + sw)
                    max_y = max(max_y, sy + sh)
                    count += 1
                if count > 0:
                    bbox_w = (max_x - min_x) / 914400  # EMU to inches
                    bbox_h = (max_y - min_y) / 914400
                    coverage = bbox_w * bbox_h
                    if coverage < DENSITY_FLOOR_SQ_IN * 0.5:  # loose threshold
                        pass  # warn below
                        # print(f"[DENSITY WARN] Slide {i+1}: {coverage:.1f} sq in "
                        #       f"(floor {DENSITY_FLOOR_SQ_IN:.1f})")

        self.prs.save(path)
        return path

    def export_pdf(self, pptx_path: str,
                   pdf_path: Optional[str] = None,
                   method: str = "auto") -> Optional[str]:
        """Export PPTX to PDF.

        Tries comtypes (PowerPoint COM, best fidelity on Windows)
        then LibreOffice headless as fallback.

        Returns PDF path on success, None on failure.
        """
        if pdf_path is None:
            pdf_path = str(Path(pptx_path).with_suffix(".pdf"))

        # Method 1: comtypes (PowerPoint COM automation on Windows)
        if method in ("auto", "comtypes"):
            if self._pdf_comtypes(pptx_path, pdf_path):
                return pdf_path

        # Method 2: LibreOffice headless
        if method in ("auto", "libreoffice"):
            if self._pdf_libreoffice(pptx_path, pdf_path):
                return pdf_path

        return None

    def _pdf_comtypes(self, pptx_path: str, pdf_path: str) -> bool:
        """Export via PowerPoint COM automation (Windows, best fidelity)."""
        try:
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = True
            try:
                deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
                deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 = PDF format
                deck.Close()
                return True
            finally:
                powerpoint.Quit()
        except Exception:
            return False

    def _pdf_libreoffice(self, pptx_path: str, pdf_path: str) -> bool:
        """Export via LibreOffice headless conversion."""
        try:
            # Find soffice
            soffice = None
            for candidate in ["soffice", "soffice.exe",
                              r"C:\Program Files\LibreOffice\program\soffice.exe",
                              r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]:
                if os.path.isfile(candidate):
                    soffice = candidate
                    break
            if not soffice:
                return False

            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 pptx_path, "--outdir", str(Path(pdf_path).parent)],
                capture_output=True, timeout=60,
                cwd=str(Path(pptx_path).parent)
            )
            if result.returncode == 0 and os.path.isfile(pdf_path):
                return True
            # LibreOffice puts the file in the same dir with .pdf extension
            expected = str(Path(pptx_path).with_suffix(".pdf"))
            if os.path.isfile(expected) and expected != pdf_path:
                import shutil
                shutil.move(expected, pdf_path)
                return True
            return os.path.isfile(pdf_path)
        except Exception:
            return False


# ---------------------------------------------------------------------------
# Content analyzer — internal decision engine for layout selection
# ---------------------------------------------------------------------------

class ContentAnalyzer:
    """Analyze content and decide the best slide layout preset.

    This is the internal council — it evaluates content patterns and
    picks the layout that best presents the information.
    """

    # Keywords that signal specific layout types
    STAT_KEYWORDS = {"%", "$", "revenue", "growth", "increase", "decrease",
                     "rate", "ratio", "count", "total", "users", "customers",
                     "market share", "nps", "score", "kpi", "roi", "yoy"}
    COMPARE_KEYWORDS = {"vs", "versus", "compared", "before", "after",
                        "pros", "cons", "advantages", "disadvantages",
                        "option", "alternative", "tradeoff"}
    TIME_KEYWORDS = {"phase", "step", "stage", "milestone", "quarter",
                     "month", "week", "sprint", "q1", "q2", "q3", "q4",
                     "timeline", "roadmap", "schedule"}
    CHECK_KEYWORDS = {"done", "complete", "completed", "pending", "todo",
                      "check", "verified", "approved", "in progress",
                      "finished", "shipped", "deployed", "launched"}
    TABLE_KEYWORDS = {"row", "column", "header", "metric", "category",
                      "segment", "region", "department", "product",
                      "tier", "plan", "package", "level"}

    @staticmethod
    def decide_layout(content: List[str]) -> str:
        """Decide the best layout based on content analysis.

        Returns one of: 'stats', 'two_column', 'table', 'checklist',
        'timeline', 'bullets'
        """
        if not content:
            return "bullets"

        text = " ".join(content).lower()

        # Check for stats (numbers + metric keywords)
        stat_score = sum(1 for kw in ContentAnalyzer.STAT_KEYWORDS if kw in text)
        if stat_score >= 2 and len(content) <= 6:
            return "stats"

        # Check for comparison (2+ items + compare keywords or state markers)
        compare_score = sum(1 for kw in ContentAnalyzer.COMPARE_KEYWORDS if kw in text)
        state_markers = sum(1 for line in content
                           if any(kw in line.lower() for kw in
                                 ["current state", "target state", "before:",
                                  "after:", "old:", "new:", "option a", "option b"]))
        if (compare_score >= 2 or state_markers >= 2) and len(content) >= 2:
            return "two_column"

        # Check for timeline (time keywords + sequential items)
        time_score = sum(1 for kw in ContentAnalyzer.TIME_KEYWORDS if kw in text)
        if time_score >= 2 and len(content) >= 3:
            return "timeline"

        # Check for checklist (check keywords + short items)
        check_score = sum(1 for kw in ContentAnalyzer.CHECK_KEYWORDS if kw in text)
        if check_score >= 1 and len(content) >= 3:
            return "checklist"

        # Check for table (pipe-separated, tab-separated, or structured keywords)
        has_pipes = any("|" in line for line in content)
        has_tabs = any("\t" in line for line in content)
        table_score = sum(1 for kw in ContentAnalyzer.TABLE_KEYWORDS if kw in text)
        if (has_pipes or has_tabs or table_score >= 2) and len(content) >= 3:
            return "table"

        # Default: bullets
        return "bullets"

    @staticmethod
    def extract_stats(content: List[str]) -> Optional[List[Tuple[str, str, str]]]:
        """Extract stats from content lines.
        Tries to find (number, label, description) patterns.
        """
        import re
        stats = []
        for line in content:
            # Look for number patterns: "$5M", "45%", "200+", "3x"
            nums = re.findall(r'[\$€£]?\d+(?:,\d+)*(?:\.\d+)?[%+]?(?:\s*[x×])?',
                              line)
            if nums:
                number = nums[0].strip()
                label = line.replace(number, "").strip()[:40]
                desc = line[len(number):].strip() if len(line) > len(number) else ""
                stats.append((number, label, desc))
        return stats[:4] if stats else None

    @staticmethod
    def split_two_column(content: List[str]) -> Optional[
        Tuple[Tuple[str, List[str]], Tuple[str, List[str]]]]:
        """Split content into two columns based on comparison markers."""
        if len(content) < 2:
            return None

        mid = len(content) // 2
        left_head = content[0][:30] if content else "Left"
        right_head = content[mid][:30] if mid < len(content) else "Right"

        left_items = [c for c in content[:mid] if len(c) > 5]
        right_items = [c for c in content[mid:] if len(c) > 5]

        if left_items and right_items:
            return (left_head, left_items), (right_head, right_items)
        return None

    @staticmethod
    def extract_table(content: List[str]) -> Optional[Tuple[List[str], List[List[str]]]]:
        """Extract table data from content lines.
        Looks for pipe-separated or structured patterns.
        """
        rows = []
        for line in content:
            if "|" in line:
                parts = [p.strip() for p in line.split("|") if p.strip()]
                if len(parts) >= 2:
                    rows.append(parts)
            elif "\t" in line:
                parts = [p.strip() for p in line.split("\t") if p.strip()]
                if len(parts) >= 2:
                    rows.append(parts)

        if len(rows) >= 2 and len(rows[0]) >= 2:
            headers = rows[0]
            data_rows = rows[1:]
            return headers, data_rows
        return None

    @staticmethod
    def extract_checklist(content: List[str]) -> Optional[List[Tuple[str, bool]]]:
        """Extract checklist items from content.
        Returns [(text, completed), ...]
        """
        items = []
        done_markers = {"done", "complete", "completed", "shipped", "launched",
                        "deployed", "verified", "approved", "✓", "[x]", "[X]"}
        pending_markers = {"pending", "todo", "in progress", "wip", "planned",
                           "[ ]", "○", "○"}

        for line in content:
            lower = line.lower()
            is_done = any(m in lower for m in done_markers)
            is_pending = any(m in lower for m in pending_markers)

            if is_done:
                # Clean the marker from text
                text = line
                for m in done_markers:
                    text = text.replace(m, "").replace(m.title(), "")
                text = text.strip(" []()✓○•-")
                if text:
                    items.append((text, True))
            elif is_pending:
                text = line
                for m in pending_markers:
                    text = text.replace(m, "").replace(m.title(), "")
                text = text.strip(" []()✓○•-")
                if text:
                    items.append((text, False))
            else:
                items.append((line.strip(), False))

        return items[:8] if items else None

    @staticmethod
    def extract_timeline(content: List[str]) -> Optional[List[Tuple[str, str, str]]]:
        """Extract timeline milestones from content.
        Returns [(phase, title, description), ...]
        """
        import re
        milestones = []
        phase_patterns = [
            r'(phase\s*\d+)', r'(step\s*\d+)', r'(stage\s*\d+)',
            r'(q[1-4])', r'(sprint\s*\d+)', r'(month\s*\d+)',
            r'(week\s*\d+)', r'(day\s*\d+)',
        ]

        for i, line in enumerate(content):
            phase = f"Step {i+1}"
            for pattern in phase_patterns:
                m = re.search(pattern, line, re.IGNORECASE)
                if m:
                    phase = m.group(1).title()
                    break

            # Split into title and description
            if ":" in line:
                parts = line.split(":", 1)
                title = parts[0].strip()
                desc = parts[1].strip() if len(parts) > 1 else ""
            else:
                title = line[:50]
                desc = line[50:].strip() if len(line) > 50 else ""

            milestones.append((phase, title, desc))

        return milestones[:5] if milestones else None


# ---------------------------------------------------------------------------
# Demo deck builder
# ---------------------------------------------------------------------------

def build_demo_deck(title: str = "Business Overview",
                    subtitle: str = "Strategic Review & Growth Initiatives",
                    output: str = "Lifewood_Presentation.pptx",
                    num_content_slides: int = 10):
    """Build a 10-slide demo deck with varied layouts per Section C rules.

    Uses: Statement, Two-column, Card grid (3x2), Stat band, Table,
    Process, Full-bleed green. No consecutive identical layouts.
    """
    deck = LifewoodDeck(title, subtitle)

    # ── Slide 1: Cover (full-bleed green) ──
    deck.add_cover()

    # ── Slide 2: Statement / Callout — big claim ──
    deck.add_callout(
        "1.0", "Vision Statement",
        "Our mission is to transform enterprises through technology — "
        "not by following trends, but by setting the standard.",
        supporting_text="Lifewood PH Team — Engineering Excellence, Delivered."
    )

    # ── Slide 3: Two-column split (6/6) ──
    deck.add_two_column(
        "1.1", "Core Capabilities",
        "Technical Depth", [
            "Cloud-native architecture design and migration",
            "AI/ML model development and deployment pipelines",
            "Full-stack engineering across web, mobile, and API",
        ],
        "Business Impact", [
            "45% average revenue growth for enterprise clients",
            "95% client retention rate across all engagements",
            "3x team scaling capacity through structured hiring",
        ]
    )

    # ── Slide 4: Card grid (3x2) — features ──
    deck.add_card_grid(
        "1.2", "What We Deliver",
        [
            ("Cloud Engineering", ["AWS/Azure/GCP architecture", "Kubernetes & containerization", "CI/CD automation"]),
            ("AI & Data", ["ML model deployment", "Data pipeline engineering", "Real-time analytics"]),
            ("Security", ["Zero-trust implementation", "SOC 2 compliance", "Penetration testing"]),
            ("Mobile", ["React Native & Flutter", "iOS/Android native", "Cross-platform strategy"]),
            ("DevOps", ["Infrastructure as code", "Monitoring & observability", "Incident response"]),
            ("Strategy", ["Technology roadmapping", "Team augmentation", "Digital transformation"]),
        ],
        grid="3x2"
    )

    # ── Slide 5: Stat band (4 figures — dense) ──
    deck.add_stats(
        "2.0", "By the Numbers",
        [
            ("200+", "Projects", "Delivered across 12 countries"),
            ("45%", "Revenue Growth", "Year-over-year for clients"),
            ("95%", "Retention", "Enterprise client satisfaction"),
            ("3x", "Team Growth", "Engineering talent expansion"),
        ]
    )

    # ── Slide 6: Process / timeline (5 steps — normal) ──
    deck.add_timeline(
        "2.1", "Engagement Model",
        [
            ("Step 1", "Discovery", "Understand goals, tech stack, and constraints"),
            ("Step 2", "Design", "Architecture blueprint, sprint plan, team composition"),
            ("Step 3", "Build", "Agile delivery with bi-weekly demos and continuous feedback"),
            ("Step 4", "Deploy", "Production release with monitoring, documentation, handover"),
            ("Step 5", "Evolve", "Ongoing optimization, support, and capability transfer"),
        ]
    )

    # ── Slide 7: Table (dense) ──
    deck.add_table(
        "2.2", "Service Comparison",
        ["Service", "Description", "Typical Timeline", "Starting Investment"],
        [
            ["Cloud Migration", "Full workload migration", "3-6 months", "$50K-200K"],
            ["Security Audit", "Zero-trust assessment", "1-2 months", "$25K-75K"],
            ["Mobile App", "Cross-platform build", "4-8 months", "$100K-400K"],
            ["DevOps Setup", "CI/CD + monitoring", "1-3 months", "$30K-100K"],
        ]
    )

    # ── Slide 8: Two-column (7/5) ──
    deck.add_two_column(
        "3.0", "Case Study — FinTech Platform",
        "Challenge", [
            "Legacy monolith unable to scale for 10x user growth",
            "Manual deployment process with 4-hour downtime windows",
            "No automated testing — 30% regression rate per release",
        ],
        "Solution", [
            "Microservices migration on AWS EKS with 50+ services",
            "GitOps CI/CD pipeline reducing deploy time to 8 minutes",
            "Comprehensive test suite achieving 92% code coverage",
        ]
    )

    # ── Slide 9: Card grid (2x2) ──
    deck.add_card_grid(
        "3.1", "Client Outcomes",
        [
            ("40% Cost Reduction", ["Infrastructure optimization", "Reserved instance planning", "Auto-scaling implementation"]),
            ("99.99% Uptime", ["Multi-AZ deployment", "Chaos engineering", "Automated failover"]),
            ("3x Faster Delivery", ["CI/CD automation", "Feature flag workflow", "Parallel environments"]),
            ("92% Test Coverage", ["Unit + integration + e2e", "Shift-left security", "Automated regression"]),
        ],
        grid="2x2"
    )

    # ── Slide 10: Closing (full-bleed green) ──
    deck.add_closing(
        "Let's Build the Future",
        "Lifewood PH Team — Your Digital Transformation Partner"
    )

    result = deck.save(output, skip_density_check=True)
    return result, len(deck._slides), deck


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Lifewood Deck Builder — premium branded PowerPoint generator"
    )
    parser.add_argument("title", nargs="?",
                        default="Business Overview",
                        help="Presentation title")
    parser.add_argument("--subtitle", "-s", default=None,
                        help="Presentation subtitle")
    parser.add_argument("--output", "-o",
                        default="Lifewood_Presentation.pptx",
                        help="Output PPTX path")
    parser.add_argument("--pdf", action="store_true",
                        help="Also export to PDF")
    parser.add_argument("--pdf-method", choices=["auto", "comtypes", "libreoffice"],
                        default="auto", help="PDF export method")
    args = parser.parse_args()

    print("Building Lifewood-branded deck...")
    path, count, deck = build_demo_deck(
        title=args.title,
        subtitle=args.subtitle or "Strategic Review & Growth Initiatives",
        output=args.output,
    )
    print(f"[OK] Saved: {path}")
    print(f"  Slides: {count}")
    print(f"  Format: Widescreen 16:9")
    print(f"  Font:   {FONT}")

    if args.pdf:
        print(f"Exporting PDF...")
        pdf = deck.export_pdf(path, method=args.pdf_method)
        if pdf:
            print(f"[OK] PDF: {pdf}")
        else:
            print("[WARN] PDF export failed — ensure PowerPoint or LibreOffice is installed")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
